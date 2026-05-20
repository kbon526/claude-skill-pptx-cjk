"""
minimal_deck.py — 最小化 5 页中文 PPT 示例。

演示本 skill 的核心 patterns：
- 不依赖外部模板，纯 python-pptx 从零创建（用于演示，实际项目优先用模板克隆）
- 三色配色锁定（#062032 深蓝黑 / #05B040 亮绿 / #FBAE40 橙）
- set_text 字体规范封装
- Hero 大字页 / 卡片 / 对比表格 / 底部 Banner

运行:
    python3 minimal_deck.py
    # 输出 minimal_deck.pptx 在当前目录

依赖：python-pptx
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    sys.stderr.write(
        "错误：未安装 python-pptx。\n"
        "macOS: /Library/Frameworks/Python.framework/Versions/3.13/bin/python3 -m pip install python-pptx\n"
    )
    sys.exit(1)

# 把 scripts 加入 sys.path 引用 set_text helper
sys.path.insert(0, str(Path(__file__).parent.parent / "scripts"))
from set_text import set_text, set_paragraphs, FONT_TITLE, FONT_BODY  # noqa: E402


# 三色配色
class Palette:
    PRIMARY = "062032"   # 深蓝黑
    SECONDARY = "05B040"  # 亮绿
    ACCENT = "FBAE40"    # 橙
    BG = "FFFFFF"
    TEXT = "333333"
    TEXT_MUTE = "8A8A8A"
    GRAY_FILL = "F5F5F5"
    GRAY_LINE = "E0E0E0"


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_full_bg(slide, hex_color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0),
                                 SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(hex_color)
    bg.line.fill.background()
    # 送到最后
    sp_tree = bg._element.getparent()
    sp_tree.remove(bg._element)
    # 找到第一个 sp 元素的位置插入
    for i, child in enumerate(sp_tree):
        if child.tag.endswith("sp"):
            sp_tree.insert(i, bg._element)
            break
    else:
        sp_tree.insert(0, bg._element)
    return bg


def add_textbox(slide, x, y, w, h, text, **kwargs):
    box = slide.shapes.add_textbox(x, y, w, h)
    set_text(box, text, **kwargs)
    return box


def add_card(slide, x, y, w, h, *, num, title, body):
    """三栏卡片之一：浅灰底 + 编号圆 + 标题 + 正文"""
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(Palette.GRAY_FILL)
    bg.line.color.rgb = RGBColor.from_string(Palette.GRAY_LINE)

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x + Inches(0.3), y + Inches(0.3),
        Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor.from_string(Palette.SECONDARY)
    circle.line.fill.background()
    set_text(circle, str(num), font=FONT_TITLE, size=18,
             bold=True, color="FFFFFF", align="center")

    title_box = slide.shapes.add_textbox(
        x + Inches(0.3), y + Inches(1.05),
        w - Inches(0.6), Inches(0.5))
    set_text(title_box, title, font=FONT_TITLE, size=16,
             bold=True, color=Palette.TEXT)

    body_box = slide.shapes.add_textbox(
        x + Inches(0.3), y + Inches(1.6),
        w - Inches(0.6), h - Inches(1.8))
    set_text(body_box, body, font=FONT_BODY, size=12,
             color=Palette.TEXT_MUTE, line_spacing=1.4)


def add_bottom_banner(slide, text):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.0),
        SLIDE_W, Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(Palette.GRAY_FILL)
    bar.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.0),
        Inches(0.06), Inches(0.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor.from_string(Palette.SECONDARY)
    accent.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(0.25), Inches(7.05),
        SLIDE_W - Inches(0.4), Inches(0.4))
    set_text(tb, text, font=FONT_BODY, size=11, color=Palette.TEXT)


def add_real_table(slide, x, y, w, h, data,
                   *, header_bg=Palette.PRIMARY,
                   header_fg="FFFFFF",
                   even_row_bg=Palette.GRAY_FILL):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = table_shape.table

    for r, row_data in enumerate(data):
        for c, text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(text)
            run.font.name = FONT_BODY
            run.font.size = Pt(12)

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(header_bg)
                run.font.color.rgb = RGBColor.from_string(header_fg)
                run.font.bold = True
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(even_row_bg)
                run.font.color.rgb = RGBColor.from_string(Palette.TEXT)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
                run.font.color.rgb = RGBColor.from_string(Palette.TEXT)
    return table_shape


# ---------- Slide builders ----------

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 通常索引 6 是空白
    add_full_bg(slide, Palette.PRIMARY)

    add_textbox(slide,
                Inches(0.8), Inches(2.5),
                Inches(11), Inches(2),
                text="2026 H2 媒介策略",
                font=FONT_TITLE, size=54, bold=True, color="FFFFFF")

    add_textbox(slide,
                Inches(0.8), Inches(4.5),
                Inches(11), Inches(0.8),
                text="精准触达 · 高效转化 · 长效增长",
                font=FONT_BODY, size=22, color="CCCCCC")

    # 右下角点缀
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(5.6), Inches(0.8), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor.from_string(Palette.SECONDARY)
    accent.line.fill.background()


def slide_section(prs, idx, title_zh):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, Palette.PRIMARY)

    add_textbox(slide,
                Inches(0.8), Inches(2.5),
                Inches(2), Inches(1.5),
                text=f"{idx:02d}",
                font=FONT_TITLE, size=72, bold=True,
                color=Palette.SECONDARY)

    add_textbox(slide,
                Inches(0.8), Inches(4.0),
                Inches(11), Inches(1),
                text=title_zh,
                font=FONT_TITLE, size=40, bold=True, color="FFFFFF")


def slide_content_cards(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    add_textbox(slide,
                Inches(0.5), Inches(0.5),
                Inches(12), Inches(0.8),
                text="三大核心策略",
                font=FONT_TITLE, size=28, bold=True, color=Palette.PRIMARY)

    # 装饰短线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.45), Inches(0.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string(Palette.SECONDARY)
    line.line.fill.background()

    # 三张卡片
    card_w = Inches(4.0)
    card_h = Inches(4.5)
    gap = Inches(0.2)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2

    cards = [
        ("01", "精准人群",
         "基于 1P + 3P 数据构建核心人群，覆盖高净值订阅家庭"),
        ("02", "全场景触达",
         "CTV + Audio + Social 三维联动，黄金时段重点投放"),
        ("03", "长效衡量",
         "品牌力 + 销量 + 客单价三维 KPI，按周复盘优化"),
    ]
    for i, (num, title, body) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, Inches(2),
                 card_w, card_h,
                 num=num, title=title, body=body)

    add_bottom_banner(slide, "通过精准媒介组合，实现品牌升级与销量突破双重目标")


def slide_comparison_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_textbox(slide,
                Inches(0.5), Inches(0.5),
                Inches(12), Inches(0.8),
                text="渠道对比",
                font=FONT_TITLE, size=28, bold=True, color=Palette.PRIMARY)

    add_real_table(slide,
        Inches(0.8), Inches(1.8),
        Inches(11.7), Inches(4.5),
        data=[
            ["维度", "CTV", "Audio", "Social"],
            ["覆盖人群", "5M+", "8M+", "12M+"],
            ["CPM", "$12", "$6", "$4"],
            ["品牌契合", "高", "中", "中"],
            ["建议占比", "45%", "25%", "30%"],
        ])

    add_bottom_banner(slide, "三渠道协同：CTV 立形象 / Audio 拓覆盖 / Social 促转化")


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, Palette.PRIMARY)

    add_textbox(slide,
                Inches(0.8), Inches(2.8),
                Inches(11.5), Inches(1.5),
                text="期待与您共创 H2 增长",
                font=FONT_TITLE, size=44, bold=True, color="FFFFFF")

    add_textbox(slide,
                Inches(0.8), Inches(4.5),
                Inches(11.5), Inches(0.6),
                text="THANK YOU",
                font=FONT_BODY, size=14, color=Palette.SECONDARY)


def main():
    prs = Presentation()
    # 16:9 默认
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_section(prs, 1, "三大核心策略")
    slide_content_cards(prs)
    slide_comparison_table(prs)
    slide_closing(prs)

    output = Path(__file__).parent / "minimal_deck.pptx"
    prs.save(output)
    print(f"已生成：{output}")
    print(f"共 {len(prs.slides)} 张幻灯片")


if __name__ == "__main__":
    main()
