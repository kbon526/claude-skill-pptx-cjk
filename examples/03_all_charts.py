"""18 种图表全集示例 — 一页一图,共 18 张 + 1 张封面 = 19 页。

运行:
    cd ~/Desktop/claude-skill-pptx-cjk
    python3.13 examples/03_all_charts.py

预期产出:
    output/All_Charts_Showcase.pptx
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx.util import Inches

from scripts.core.color import (
    BRAND_COLORS, register_brand, RGBColor,
    RED, BLUE, GREEN, ORANGE, BLACK,
    CLR_FACEBOOK, CLR_INSTAGRAM, CLR_TIKTOK, CLR_YOUTUBE,
    CLR_XHS, CLR_DOUYIN, CLR_WECHAT,
    PHASE_AWARENESS, PHASE_INTEREST, PHASE_CONVERSION,
    LIGHT_GRAY, MED_GRAY,
)
from scripts.core.font import patch_theme_fonts
from scripts.core.registry import (
    init_registry, add_slide, attach_visual, save,
    set_checkpoint, CKPT_APPROVED,
)
from scripts.pipeline import chart_remake as cr
from scripts.pipeline.assemble import assemble
from scripts.components import hero as comp_hero

OUTPUT = Path("output/All_Charts_Showcase.pptx")
TEMPLATE = Path("templates/starter.pptx")
OUTPUT.parent.mkdir(parents=True, exist_ok=True)

if not TEMPLATE.exists():
    print("⚠️ Template missing. Run: python3 scripts/tools/generate_starter_template.py templates/starter.pptx")
    sys.exit(1)


# ══════════════════════════════════════════════════════════
# 准备 registry
# ══════════════════════════════════════════════════════════
reg = init_registry("18 种图表展示", brand="Showcase", version="v1")

# Cover
add_slide(reg, "slide_01",
          title="18 种图表类型展示",
          subtitle="A 类原生 12 种(可双击编辑) + B 类形状 7 种",
          layout="section_divider")


# ══════════════════════════════════════════════════════════
# A 类 — 12 个原生图表(双击可编辑数据)
# ══════════════════════════════════════════════════════════

# 1. 柱状图
add_slide(reg, "slide_02",
          title="1. 柱状图(column)",
          subtitle="单系列 — 不同实体的同一指标对比",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_02",
    cr.make_column_spec(
        cats=["品牌 A", "品牌 B", "品牌 C", "品牌 D"],
        vals=[38, 23, 24, 15],
        colors=[RED, BLUE, GREEN, ORANGE],
        fmt='0"%"',
    ))

# 2. 条形图
add_slide(reg, "slide_03",
          title="2. 条形图(bar)",
          subtitle="水平条形,适合长名标签",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_03",
    cr.make_bar_spec(
        cats=["品牌资产", "市场地位", "产品力", "数字营销", "外部环境"],
        vals=[3.18, 3.08, 3.35, 3.88, 4.13],
        colors=[BLUE] * 5,
        fmt='0.00',
    ))

# 3. 分组柱状图
add_slide(reg, "slide_04",
          title="3. 分组柱状图(grouped_col)",
          subtitle="多系列 × 多类目,同一坐标系对比",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_04",
    cr.make_grouped_col_spec(
        cats=["Q1", "Q2", "Q3", "Q4"],
        series={
            "品牌 A": ([100, 120, 140, 130], RED),
            "品牌 B": ([80, 90, 100, 110], BLUE),
            "AcmeCorp": ([60, 75, 85, 95], GREEN),
        },
    ))

# 4. 堆叠柱状图(100%)
add_slide(reg, "slide_05",
          title="4. 堆叠柱状图(stacked_col 100%)",
          subtitle="多系列占比变化(时间维度)",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_05",
    cr.make_stacked_col_spec(
        cats=["1月", "2月", "3月", "4月"],
        series_list=[
            ("Facebook", [40, 38, 35, 32], CLR_FACEBOOK),
            ("Instagram", [25, 27, 30, 33], CLR_INSTAGRAM),
            ("TikTok", [20, 22, 23, 25], CLR_TIKTOK),
            ("YouTube", [15, 13, 12, 10], CLR_YOUTUBE),
        ],
        percent=True,
    ))

# 5. 堆叠条形图(100%)
add_slide(reg, "slide_06",
          title="5. 堆叠条形图(stacked_bar 100%)",
          subtitle="跨实体的渠道分布对比",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_06",
    cr.make_stacked_bar_spec(
        cats=["品牌 A", "品牌 B", "AcmeCorp"],
        series_list=[
            ("小红书", [22, 35, 28], CLR_XHS),
            ("抖音", [42, 25, 32], CLR_DOUYIN),
            ("微信", [18, 25, 22], CLR_WECHAT),
            ("其他", [18, 15, 18], MED_GRAY),
        ],
        percent=True,
    ))

# 6. 饼图
add_slide(reg, "slide_07",
          title="6. 饼图(pie)",
          subtitle="单层占比,适合 ≤5 类",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_07",
    cr.make_pie_spec(
        cats=["A", "B", "C", "D"],
        vals=[40, 30, 20, 10],
        colors=[RED, BLUE, GREEN, ORANGE],
    ))

# 7. 环形图
add_slide(reg, "slide_08",
          title="7. 环形图(doughnut)",
          subtitle="比饼图更现代,中心可放主指标",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_08",
    cr.channel_mix_spec({
        "Facebook": 43, "Instagram": 23,
        "TikTok": 19, "YouTube": 15,
    }))

# 8. 折线图
add_slide(reg, "slide_09",
          title="8. 折线图(line)",
          subtitle="时间序列趋势,多系列对比",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_09",
    cr.make_line_spec(
        cats=["1月", "2月", "3月", "4月", "5月", "6月"],
        series={
            "GMV": ([100, 120, 140, 130, 160, 180], RED),
            "订单": ([80, 95, 110, 105, 130, 145], BLUE),
        },
        smooth=False, markers=True,
    ))

# 9. 面积图
add_slide(reg, "slide_10",
          title="9. 面积图(area / stacked)",
          subtitle="趋势 + 累计量,堆叠版可显示构成变化",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_10",
    cr.make_area_spec(
        cats=["1月", "2月", "3月", "4月", "5月"],
        series={
            "新客": ([20, 25, 30, 35, 40], RED),
            "老客": ([60, 65, 70, 80, 90], BLUE),
            "唤回": ([10, 12, 15, 18, 20], GREEN),
        },
        stacked=True,
    ))

# 10. 主题河流图
add_slide(reg, "slide_11",
          title="10. 主题河流图(theme_river)",
          subtitle="多系列累计趋势,中心对称版需 PPT 后期手调",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_11",
    cr.make_theme_river_spec(
        cats=["Q1", "Q2", "Q3", "Q4"],
        series={
            "话题 A": ([10, 25, 40, 30], RED),
            "话题 B": ([20, 30, 35, 45], BLUE),
            "话题 C": ([15, 20, 25, 50], GREEN),
        },
    ))

# 11. 雷达图
add_slide(reg, "slide_12",
          title="11. 雷达图(radar)",
          subtitle="多维评分对比",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_12",
    cr.make_radar_spec(
        cats=["品牌资产", "市场地位", "产品力", "数字营销", "外部环境"],
        series={
            "品牌 A": ([5.0, 3.7, 4.5, 5.0, 3.4], RED),
            "AcmeCorp": ([3.2, 3.1, 3.4, 3.9, 4.1], BLUE),
        },
    ))

# 12. 散点图
add_slide(reg, "slide_13",
          title="12. 散点图(scatter)",
          subtitle="二变量相关性",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_13",
    cr.make_scatter_spec(
        series={
            "组 A": [(1, 2), (2, 3), (3, 5), (4, 7), (5, 9), (6, 10)],
            "组 B": [(1, 5), (2, 6), (3, 6), (4, 8), (5, 8), (6, 9)],
        },
        colors=[RED, BLUE],
    ))

# 13. 气泡图
add_slide(reg, "slide_14",
          title="13. 气泡图(bubble)",
          subtitle="三维数据(x, y, 大小)",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_14",
    cr.make_bubble_spec(
        series={
            "产品矩阵": [
                (10, 20, 30), (20, 40, 50), (30, 35, 80),
                (40, 50, 40), (50, 60, 90), (60, 45, 60),
            ],
        },
        colors=[BLUE],
    ))


# ══════════════════════════════════════════════════════════
# B 类 — 7 种 diagrams(形状组合,不可双击编辑数据)
# ══════════════════════════════════════════════════════════

# 14. K 线图(用形状画)
add_slide(reg, "slide_15",
          title="14. K 线图(candlestick)",
          subtitle="股票/金融数据 — 用形状绘制(python-pptx 不支持 STOCK_OHLC 写入)",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_15",
    cr.make_stock_spec(
        dates=["1日", "2日", "3日", "4日", "5日", "6日", "7日"],
        ohlc=[
            (100, 110, 95, 105),
            (105, 115, 100, 112),
            (112, 120, 108, 110),  # 跌
            (110, 118, 105, 115),
            (115, 125, 112, 122),
            (122, 128, 118, 120),  # 跌
            (120, 130, 117, 128),
        ],
    ))

# 15. 漏斗图
add_slide(reg, "slide_16",
          title="15. 漏斗图(funnel)",
          subtitle="转化路径单向衰减",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_16",
    cr.make_funnel_spec(
        stages=["曝光", "点击", "加购", "下单", "复购"],
        values=[100000, 30000, 8000, 2000, 500],
        colors=[
            RGBColor(0xE6, 0x00, 0x00),
            RGBColor(0xE6, 0x40, 0x30),
            RGBColor(0xCC, 0x70, 0x60),
            RGBColor(0xB0, 0xA0, 0x90),
            RGBColor(0x99, 0xCC, 0xCC),
        ],
    ))

# 16. 流程图
add_slide(reg, "slide_17",
          title="16. 流程图(process_flow)",
          subtitle="阶段性步骤,横向箭头连接",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_17",
    cr.make_process_flow_spec(
        steps=[
            ("立势期", "Q3 7-8月"),
            ("冲刺期", "Q4 9-10月"),
            ("收割期", "Q4 11-12月"),
        ],
        colors=[PHASE_AWARENESS, PHASE_INTEREST, PHASE_CONVERSION],
    ))

# 17. 桑基图
add_slide(reg, "slide_18",
          title="17. 桑基图(sankey)",
          subtitle="左右两层流向 — 简化版",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_18",
    cr.make_sankey_spec(
        left_nodes=[("源 A", 60), ("源 B", 40)],
        right_nodes=[("目标 X", 35), ("目标 Y", 30), ("目标 Z", 35)],
        flows=[
            (0, 0, 25), (0, 1, 20), (0, 2, 15),
            (1, 0, 10), (1, 1, 10), (1, 2, 20),
        ],
        left_colors=[RED, BLUE],
        right_colors=[GREEN, ORANGE, BLACK],
    ))

# 18. 矩形树图
add_slide(reg, "slide_19",
          title="18. 矩形树图(treemap)",
          subtitle="按面积比例展示占比",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_19",
    cr.make_treemap_spec(
        items=[
            ("品类 A", 100), ("品类 B", 70), ("品类 C", 50),
            ("品类 D", 30), ("品类 E", 20), ("品类 F", 15), ("品类 G", 10),
        ],
        colors=[RED, BLUE, GREEN, ORANGE, BLACK, MED_GRAY,
                RGBColor(0x80, 0x80, 0xFF)],
    ))

# 19. 旭日图(简化双层环形)
add_slide(reg, "slide_20",
          title="19. 旭日图(sunburst,简化双层环形)",
          subtitle="多层级占比 — python-pptx 限制,只做内圈环形示意",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_20",
    cr.make_sunburst_spec(
        hierarchy={
            "线上": {"小红书": 25, "抖音": 30, "微信": 15},
            "线下": {"门店": 20, "活动": 10},
        },
        colors={"线上": BLUE, "线下": ORANGE},
    ))

# 20. 甘特图
add_slide(reg, "slide_21",
          title="20. 甘特图(gantt)",
          subtitle="任务排程 — 用 BAR_STACKED 第一系列透明做偏移",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_21",
    cr.make_gantt_spec(
        tasks=[
            ("需求调研", 0, 5),
            ("方案撰写", 4, 8),
            ("内部 review", 11, 3),
            ("客户呈现", 14, 2),
            ("修改打磨", 16, 6),
            ("最终交付", 22, 2),
        ],
        color=BLUE,
    ))

# 21. 热力图
add_slide(reg, "slide_22",
          title="21. 热力图(heatmap)",
          subtitle="二维稠密分布 — 表格 + 渐变填充",
          layout="chart_focus")
cr.attach_chart_to_slide(reg, "slide_22",
    cr.make_heatmap_spec(
        matrix=[
            [10, 25, 40, 30, 15],
            [20, 35, 60, 45, 25],
            [15, 30, 50, 55, 35],
            [25, 40, 70, 65, 45],
            [10, 20, 35, 30, 20],
        ],
        x_labels=["周一", "周二", "周三", "周四", "周五"],
        y_labels=["8 点", "10 点", "12 点", "14 点", "16 点"],
        cmap_low=RGBColor(0xF0, 0xF0, 0xF0),
        cmap_high=RGBColor(0xE6, 0x00, 0x00),
    ))


# ══════════════════════════════════════════════════════════
# 组装(跳过 checkpoint 校验,演示模式)
# ══════════════════════════════════════════════════════════
for ck in ["ckpt_1_framework", "ckpt_2_content", "ckpt_3_visual", "ckpt_4_assembly"]:
    set_checkpoint(reg, ck, CKPT_APPROVED)

save(reg, "work/showcase_registry.json")

print("=" * 70)
print(f"开始组装 {len(reg['slides'])} 张 slide → {OUTPUT}")
print("=" * 70)

assemble(reg, template_path=TEMPLATE, output_path=OUTPUT)

print(f"\n✅ 输出: {OUTPUT}")
print(f"   含 {len(reg['slides'])} 页(1 封面 + 21 张图表演示)")
print()
print("视觉 QA:")
print(f"  python3 scripts/render_layout.py {OUTPUT} output/showcase_qa/")
