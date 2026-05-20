"""
insert_slide_at.py — python-pptx 任意位置插入 slide 的 helper。

python-pptx 的 add_slide() 只能追加到末尾。本 helper 通过操作内部 XML
(_sldIdLst) 实现任意位置插入。

用法 1（作为模块）：
    from insert_slide_at import insert_slide_at
    new_slide = insert_slide_at(prs, idx=5, layout=prs.slide_layouts[1])

用法 2（命令行验证）：
    python3 insert_slide_at.py input.pptx 5 1 output.pptx
    # 在 input.pptx 的第 5 个位置插入一张使用 layout[1] 的空白 slide
"""

import sys


def insert_slide_at(prs, idx, layout):
    """
    在 prs.slides 的指定 idx 位置插入新 slide。

    参数:
        prs: pptx.Presentation 对象
        idx: 目标位置（0-based）
        layout: prs.slide_layouts 中的某个 layout 对象

    返回:
        新创建的 slide 对象
    """
    new_slide = prs.slides.add_slide(layout)
    sld_lst = prs.slides._sldIdLst
    new_sld_id = list(sld_lst)[-1]
    sld_lst.remove(new_sld_id)

    # 把它插到 idx 位置
    if idx >= len(list(sld_lst)):
        sld_lst.append(new_sld_id)
    else:
        target = list(sld_lst)[idx]
        target.addprevious(new_sld_id)

    return new_slide


def remove_slide_at(prs, idx):
    """删除 prs.slides[idx]。"""
    from pptx.oxml.ns import qn
    sld_lst = prs.slides._sldIdLst
    sld_id_elem = list(sld_lst)[idx]
    rId = sld_id_elem.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sld_lst.remove(sld_id_elem)


def main():
    if len(sys.argv) != 5:
        sys.stderr.write(
            "用法: python3 insert_slide_at.py <input.pptx> "
            "<insert_idx> <layout_idx> <output.pptx>\n"
        )
        sys.exit(1)

    input_path, idx_str, layout_idx_str, output_path = sys.argv[1:]
    idx = int(idx_str)
    layout_idx = int(layout_idx_str)

    from pptx import Presentation
    prs = Presentation(input_path)
    layout = prs.slide_layouts[layout_idx]
    insert_slide_at(prs, idx, layout)
    prs.save(output_path)

    print(f"已在 {input_path} 的 idx={idx} 位置插入 layout[{layout_idx}] 空白 slide")
    print(f"输出: {output_path}")


if __name__ == "__main__":
    main()
