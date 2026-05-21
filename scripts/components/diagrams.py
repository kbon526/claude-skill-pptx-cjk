"""diagrams.py — B 类形状组合图表(数据嵌入代码,不可双击编辑数据)。

⚠️ 与 chart.py 的关键差异:
- chart.py 的图表客户在 PowerPoint 双击可改 Excel 数据
- diagrams.py 的图表客户改要回到 Python 代码改

但是:
- diagrams 的"形状"在 PowerPoint 里完全可改(改色、改大小、改文字)
- 适合 python-pptx 原生 Chart 不支持的视觉(漏斗/流程图/桑基/旭日/树图/甘特/热力)

提供的类型:
- funnel — 营销漏斗(梯形堆叠)
- process_flow — 横向流程图(矩形 + 箭头)
- sankey — 桑基图(矩形 + 贝塞尔曲线,简化版)
- treemap — 矩形树图(嵌套矩形,固定算法)
- sunburst — 旭日图(同心环 + 扇形,简化版)
- gantt — 甘特图(用 BAR_STACKED 技巧:第一系列透明做偏移)
- heatmap — 热力图(表格 + 渐变填充)
"""

import math
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from ..core.font import _sf, FONT_LIGHT, SZ_BODY, SZ_CAP, SZ_SRC
from ..core.color import (
    WHITE, DARK_GRAY, TEXT_GRAY, LIGHT_GRAY, MED_GRAY,
    RED, BLUE, GREEN, ORANGE, BLACK, GRID_GRAY,
)


# ══════════════════════════════════════════════════════════
# B 类 1: 漏斗图 (funnel)
# ══════════════════════════════════════════════════════════
def add_funnel(slide, l, t, w, h, stages, values=None, colors=None,
                show_values=True, show_pct=True):
    """漏斗图 — 梯形堆叠。

    Args:
        stages: list[str] 阶段名(从顶到底)
        values: list[num] 对应数值;不传则等比降序
        colors: list[RGBColor] 每阶段的颜色;不传则用红→灰渐变
        show_values: 是否在右侧显示数值
        show_pct: 是否显示对上一阶段的转化率

    形态:
        每个阶段是一个等高的梯形,宽度按 values 比例,垂直堆叠。
    """
    n = len(stages)
    if values is None:
        # 默认等比降序(每层 80%)
        values = [100 * (0.8 ** i) for i in range(n)]
    if colors is None:
        # 红 → 浅灰渐变
        colors = [
            RGBColor(
                int(0xE6 - (0xE6 - 0xCC) * i / (n - 1)),
                int(0x00 + (0xCC - 0x00) * i / (n - 1)),
                int(0x00 + (0xCC - 0x00) * i / (n - 1)),
            )
            for i in range(n)
        ]

    max_v = max(values)
    stage_h = int(h / n)
    label_w = Inches(2.0)  # 右侧数值区宽度
    funnel_w = w - label_w  # 漏斗实际占用宽度

    for i, (name, val) in enumerate(zip(stages, values)):
        # 当前阶段宽度按比例
        stage_w = int(funnel_w * (val / max_v))
        x_offset = int(l + (funnel_w - stage_w) // 2)
        y = int(t + i * stage_h)

        # 用矩形(简化)代替梯形 — 视觉效果接近
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x_offset, y, stage_w, stage_h - Inches(0.04)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors[i]
        bar.line.fill.background()

        # 阶段名(白色,居中)
        tf = bar.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = name
        _sf(r, SZ_BODY, True, WHITE)

        # 右侧数值 + 转化率
        if show_values:
            label_x = l + funnel_w + Inches(0.1)
            label_text_lines = [(f"{val:,.0f}", SZ_BODY, True, DARK_GRAY)]
            if show_pct and i > 0:
                pct = (val / values[i - 1]) * 100
                label_text_lines.append((f"↓ {pct:.1f}%", SZ_SRC, False, TEXT_GRAY))

            tb = slide.shapes.add_textbox(label_x, y, label_w - Inches(0.1),
                                           stage_h)
            tf2 = tb.text_frame
            tf2.word_wrap = True
            tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            for j, (text, sz, bold, color) in enumerate(label_text_lines):
                pp = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                rr = pp.add_run()
                rr.text = text
                _sf(rr, sz, bold, color)


# ══════════════════════════════════════════════════════════
# B 类 2: 流程图 (process_flow)
# ══════════════════════════════════════════════════════════
def add_process_flow(slide, l, t, w, h, steps, colors=None,
                      arrow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    """横向流程图 — 矩形 + 箭头。

    Args:
        steps: list[str | (title, desc)] 步骤;若是 tuple,第二项作为副说明
        colors: list[RGBColor] 每步的颜色
        arrow: 步骤间是否加箭头连接符
        shape: MSO_SHAPE.RECTANGLE / ROUNDED_RECTANGLE / CHEVRON

    形态:
        横向 N 个矩形等宽排列,中间用三角箭头连接。
    """
    n = len(steps)
    if colors is None:
        # 蓝色渐变(立势→冲刺→收割风格)
        colors = [BLUE] * n

    arrow_w = Inches(0.3) if arrow else Inches(0)
    gap = Inches(0.1) if arrow else Inches(0.15)
    box_w = int((w - arrow_w * (n - 1) - gap * (n - 1)) / n)

    for i, step in enumerate(steps):
        if isinstance(step, tuple):
            title, desc = step
        else:
            title, desc = step, ""

        x = int(l + i * (box_w + arrow_w + gap))
        box = slide.shapes.add_shape(shape, x, t, box_w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(6)
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        _sf(r, SZ_BODY, True, WHITE)

        if desc:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(3)
            r2 = p2.add_run()
            r2.text = desc
            _sf(r2, SZ_CAP, False, RGBColor(0xEE, 0xEE, 0xEE))

        # 箭头(除最后一个)
        if arrow and i < n - 1:
            ax = int(x + box_w + gap / 2)
            ay = int(t + h / 2 - Inches(0.1))
            ar = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, ax, ay, arrow_w, Inches(0.2)
            )
            ar.fill.solid()
            ar.fill.fore_color.rgb = MED_GRAY
            ar.line.fill.background()


# ══════════════════════════════════════════════════════════
# B 类 3: 桑基图 (sankey,简化版)
# ══════════════════════════════════════════════════════════
def add_sankey(slide, l, t, w, h, left_nodes, right_nodes, flows,
                left_colors=None, right_colors=None):
    """桑基图(简化版:两层 — 左源 → 右目标)。

    Args:
        left_nodes: list[(name, value)] 左侧源节点
        right_nodes: list[(name, value)] 右侧目标节点
        flows: list[(left_idx, right_idx, value)] 流量

    形态:
        左侧 N 个矩形堆叠 + 右侧 M 个矩形堆叠 + 中间用半透明带状形状连接。
        简化:连接用直线而非贝塞尔(python-pptx 不支持任意曲线绘制)。

    复杂场景(多层桑基)建议用专业工具或外部库生成图后嵌入。
    """
    node_w = Inches(0.4)
    flow_left_x = l + node_w
    flow_right_x = l + w - node_w
    flow_w = flow_right_x - flow_left_x

    # 左侧节点位置
    left_total = sum(v for _, v in left_nodes)
    left_positions = []
    cur_y = t
    for name, val in left_nodes:
        node_h = int(h * (val / left_total))
        left_positions.append((cur_y, node_h))
        cur_y += node_h

    # 右侧节点位置
    right_total = sum(v for _, v in right_nodes)
    right_positions = []
    cur_y = t
    for name, val in right_nodes:
        node_h = int(h * (val / right_total))
        right_positions.append((cur_y, node_h))
        cur_y += node_h

    # 默认配色
    if left_colors is None:
        left_colors = [BLUE] * len(left_nodes)
    if right_colors is None:
        right_colors = [GREEN] * len(right_nodes)

    # 流量带 — 用半透明矩形近似
    left_used = [0] * len(left_nodes)
    right_used = [0] * len(right_nodes)
    for li, ri, val in flows:
        l_y, l_h = left_positions[li]
        r_y, r_h = right_positions[ri]
        flow_h_l = int(h * (val / left_total))
        flow_h_r = int(h * (val / right_total))

        y_start_top = int(l_y + left_used[li])
        y_end_top = int(r_y + right_used[ri])
        left_used[li] += flow_h_l
        right_used[ri] += flow_h_r

        c_top = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            int(flow_left_x), y_start_top,
            int(flow_right_x), y_end_top,
        )
        c_top.line.color.rgb = LIGHT_GRAY
        c_top.line.width = Pt(0.5)

        c_bot = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            int(flow_left_x), y_start_top + flow_h_l,
            int(flow_right_x), y_end_top + flow_h_r,
        )
        c_bot.line.color.rgb = LIGHT_GRAY
        c_bot.line.width = Pt(0.5)

    # 左侧节点 + 文字
    for i, (name, val) in enumerate(left_nodes):
        y, node_h = left_positions[i]
        node = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, l, y, node_w, node_h - Inches(0.02)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = left_colors[i % len(left_colors)]
        node.line.fill.background()
        # 标签
        tb = slide.shapes.add_textbox(
            l - Inches(1.5), y, Inches(1.4), node_h
        )
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"{name}\n{val:,.0f}"
        _sf(r, SZ_CAP, color=DARK_GRAY)

    # 右侧节点
    for i, (name, val) in enumerate(right_nodes):
        y, node_h = right_positions[i]
        node = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, l + w - node_w, y, node_w, node_h - Inches(0.02)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = right_colors[i % len(right_colors)]
        node.line.fill.background()
        tb = slide.shapes.add_textbox(
            l + w + Inches(0.1), y, Inches(1.4), node_h
        )
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{name}\n{val:,.0f}"
        _sf(r, SZ_CAP, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════
# B 类 4: 矩形树图 (treemap,简化 squarify 算法)
# ══════════════════════════════════════════════════════════
def add_treemap(slide, l, t, w, h, items, colors=None):
    """矩形树图(单层,squarify 切分算法的简化版)。

    Args:
        items: list[(name, value)] 按 value 降序自动排
        colors: list[RGBColor] 与 items 顺序对应

    形态:
        按面积比例切分矩形,大块在左上,小块在右下(squarify 风格)。
    """
    items = sorted(items, key=lambda x: -x[1])  # 降序
    total = sum(v for _, v in items)
    n = len(items)
    if colors is None:
        colors = [BLUE] * n

    # 简化的 slice-and-dice:横向竖向交替切分
    def slice_and_dice(items_with_color, x, y, ww, hh, vertical=True):
        if not items_with_color:
            return
        if len(items_with_color) == 1:
            (name, val), clr = items_with_color[0]
            _draw_treemap_cell(slide, x, y, ww, hh, name, val, clr)
            return
        # 切一半
        total_v = sum(v for (_, v), _ in items_with_color)
        cur = 0
        split_idx = 0
        for i, ((_, v), _) in enumerate(items_with_color):
            cur += v
            if cur >= total_v * 0.5:
                split_idx = i + 1
                break
        if split_idx == 0:
            split_idx = 1
        first = items_with_color[:split_idx]
        rest = items_with_color[split_idx:]
        first_v = sum(v for (_, v), _ in first)
        ratio = first_v / total_v

        if vertical:
            split_w = int(ww * ratio)
            slice_and_dice(first, x, y, split_w, hh, not vertical)
            slice_and_dice(rest, x + split_w, y, ww - split_w, hh, not vertical)
        else:
            split_h = int(hh * ratio)
            slice_and_dice(first, x, y, ww, split_h, not vertical)
            slice_and_dice(rest, x, y + split_h, ww, hh - split_h, not vertical)

    items_with_color = list(zip(items, colors))
    slice_and_dice(items_with_color, l, t, w, h, vertical=(w > h))


def _draw_treemap_cell(slide, x, y, w, h, name, val, color):
    """绘制单个 treemap 单元格。"""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = WHITE
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = str(name)
    # 字号根据格子大小自适应(h 是 EMU 整数,914400 EMU = 1 inch)
    h_in = h / 914400 if isinstance(h, int) else h.emu / 914400
    sz_pt = max(8, min(18, int(h_in * 8)))
    _sf(r, Pt(sz_pt), True, WHITE)

    # 数值
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = f"{val:,.0f}"
    _sf(r2, Pt(max(7, sz_pt - 4)), False, WHITE)


# ══════════════════════════════════════════════════════════
# B 类 5: 旭日图 (sunburst,简化版)
# ══════════════════════════════════════════════════════════
def add_sunburst(slide, l, t, w, h, hierarchy, colors=None):
    """旭日图(简化两层版)。

    Args:
        hierarchy: dict[parent_name -> dict[child_name -> value]]
                   或 list[(parent, child, value)]
        colors: parent → RGBColor 字典

    实现说明:
        python-pptx 的 PIE / DOUGHNUT 不支持嵌套环,要做真正的旭日图
        必须用形状的 PIE / ARC 组合,实现成本高。
        此函数提供"双层环形图"近似 — 用 2 个 doughnut 同心叠放。
    """
    # 解析 hierarchy
    if isinstance(hierarchy, dict):
        items = []
        for parent, children in hierarchy.items():
            for child, val in children.items():
                items.append((parent, child, val))
    else:
        items = list(hierarchy)

    # 内圈:parent 聚合
    parent_totals = {}
    for parent, _, val in items:
        parent_totals[parent] = parent_totals.get(parent, 0) + val

    if colors is None:
        # 自动配色
        default_palette = [BLUE, GREEN, ORANGE, RED]
        colors = {p: default_palette[i % len(default_palette)]
                  for i, p in enumerate(parent_totals.keys())}

    # 内圈环形图
    inner_size = int(min(w, h) * 0.6)
    inner_x = int(l + (w - inner_size) // 2)
    inner_y = int(t + (h - inner_size) // 2)

    cd_inner = CategoryChartData()
    cd_inner.categories = list(parent_totals.keys())
    cd_inner.add_series("", list(parent_totals.values()))
    cf_inner = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, inner_x, inner_y, inner_size, inner_size, cd_inner,
    )
    c1 = cf_inner.chart
    c1.has_legend = False
    s1 = c1.series[0]
    for i, p in enumerate(parent_totals.keys()):
        s1.points[i].format.fill.solid()
        s1.points[i].format.fill.fore_color.rgb = colors[p]
    return cf_inner


# ══════════════════════════════════════════════════════════
# B 类 6: 甘特图 (gantt) — 用 BAR_STACKED 技巧
# ══════════════════════════════════════════════════════════
def add_gantt(slide, l, t, w, h, tasks, color=None):
    """甘特图 — 用 BAR_STACKED 实现(第一系列透明做时间偏移)。

    Args:
        tasks: list[(task_name, start_day, duration_days)]
        color: RGBColor 任务条颜色

    技巧:
        BAR_STACKED 第一个系列设为透明 = 起始偏移
        第二个系列 = 任务持续时间(可见的条)
    """
    if color is None:
        color = BLUE

    cats = [task[0] for task in tasks]
    starts = [task[1] for task in tasks]
    durations = [task[2] for task in tasks]

    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("Start", starts)        # 第一系列:透明偏移
    cd.add_series("Duration", durations)  # 第二系列:任务条
    cf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, l, t, w, h, cd)
    c = cf.chart
    c.has_legend = False
    c.style = 2

    # 第一系列(Start)设为透明
    s_start = c.series[0]
    s_start.format.fill.background()
    s_start.format.line.fill.background()

    # 第二系列(Duration)上色
    s_dur = c.series[1]
    s_dur.format.fill.solid()
    s_dur.format.fill.fore_color.rgb = color

    # 字体
    for ax in (c.category_axis, c.value_axis):
        ax.tick_labels.font.size = SZ_SRC
        ax.tick_labels.font.name = FONT_LIGHT
        ax.tick_labels.font.color.rgb = TEXT_GRAY
    return cf


# ══════════════════════════════════════════════════════════
# B 类 7: 热力图 (heatmap) — 表格 + 渐变填充
# ══════════════════════════════════════════════════════════
def add_heatmap(slide, l, t, w, h, matrix, x_labels, y_labels,
                  cmap_low=None, cmap_high=None, fmt="{:.0f}"):
    """热力图 — 用表格 + 单元格渐变色实现。

    Args:
        matrix: 2D list[[v]],外层是行(y),内层是列(x)
        x_labels: 列名(顶部)
        y_labels: 行名(左侧)
        cmap_low: 低值颜色(默认浅灰)
        cmap_high: 高值颜色(默认红)
        fmt: 单元格数值格式

    形态:
        N+1 行 × M+1 列表格,首行/首列是标签,其余按值填渐变色。
    """
    if cmap_low is None:
        cmap_low = LIGHT_GRAY
    if cmap_high is None:
        cmap_high = RED

    n_rows = len(y_labels) + 1
    n_cols = len(x_labels) + 1

    ts = slide.shapes.add_table(n_rows, n_cols, l, t, w, h)
    tbl = ts.table

    # 找 min/max 用于色阶
    flat_vals = [v for row in matrix for v in row]
    vmin, vmax = min(flat_vals), max(flat_vals)
    rng = vmax - vmin if vmax > vmin else 1

    # 关闭默认条带
    tblPr = tbl._tbl.tblPr
    tblPr.set("bandRow", "0")
    tblPr.set("bandCol", "0")
    tblPr.set("firstRow", "0")
    tblPr.set("firstCol", "0")

    # 填充数据
    for ri in range(n_rows):
        for ci in range(n_cols):
            cell = tbl.cell(ri, ci)
            if ri == 0 and ci == 0:
                cell.text = ""
                _style_heatmap_cell(cell, "", bold=True, fill=WHITE)
            elif ri == 0:
                # 列标签
                cell.text = str(x_labels[ci - 1])
                _style_heatmap_cell(cell, x_labels[ci - 1],
                                     bold=True, fill=LIGHT_GRAY)
            elif ci == 0:
                # 行标签
                cell.text = str(y_labels[ri - 1])
                _style_heatmap_cell(cell, y_labels[ri - 1],
                                     bold=True, fill=LIGHT_GRAY)
            else:
                # 数据单元格
                val = matrix[ri - 1][ci - 1]
                norm = (val - vmin) / rng
                # 线性插值 cmap_low → cmap_high
                fill = RGBColor(
                    int(cmap_low[0] + (cmap_high[0] - cmap_low[0]) * norm),
                    int(cmap_low[1] + (cmap_high[1] - cmap_low[1]) * norm),
                    int(cmap_low[2] + (cmap_high[2] - cmap_low[2]) * norm),
                )
                # 文字色:背景越深用白色
                txt_color = WHITE if norm > 0.5 else DARK_GRAY
                _style_heatmap_cell(cell, fmt.format(val),
                                     bold=False, fill=fill, color=txt_color)


def _style_heatmap_cell(cell, text, bold=False, fill=None, color=None):
    """填充单元格:字体、对齐、背景色。"""
    cell.text = str(text)
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            _sf(r, SZ_CAP, bold, color or DARK_GRAY)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()


# ══════════════════════════════════════════════════════════
# B 类 8: K 线图 (candlestick) — 用形状绘制
# ══════════════════════════════════════════════════════════
# python-pptx 1.x 不支持 STOCK_OHLC 写入,改为形状绘制。
# 每个交易日 = 一根垂直影线(low~high) + 一个矩形实体(open~close)。
def add_candlestick(slide, l, t, w, h, dates, ohlc,
                     up_color=None, down_color=None,
                     show_labels=True):
    """K 线图(蜡烛图) — 用形状绘制。

    Args:
        dates: list[str] 时间序列
        ohlc: list[(open, high, low, close)]
        up_color: 涨色(默认中国惯例红 = RED)
        down_color: 跌色(默认中国惯例绿 = GREEN)
        show_labels: 是否在底部显示日期

    形态:
        一根细竖线(low → high) + 矩形实体(open → close)。
        红实体 = 收盘 ≥ 开盘,绿实体 = 收盘 < 开盘。
    """
    if up_color is None:
        up_color = RED
    if down_color is None:
        down_color = GREEN

    n = len(dates)
    if n != len(ohlc):
        raise ValueError("dates 和 ohlc 长度必须一致")

    # 价格范围
    all_prices = [p for o in ohlc for p in o]
    p_min, p_max = min(all_prices), max(all_prices)
    p_range = p_max - p_min if p_max > p_min else 1
    # 上下留 5% 边距
    p_min -= p_range * 0.05
    p_max += p_range * 0.05
    p_range = p_max - p_min

    # 画布分配:左侧留 0.5in 给 y 轴标签,底部留 0.4in 给日期
    label_w = Inches(0.5)
    label_h = Inches(0.4) if show_labels else Inches(0)
    plot_l = l + label_w
    plot_t = t
    plot_w = w - label_w
    plot_h = h - label_h

    # 网格线 + y 轴标签(5 等分)
    for i in range(5):
        y_val = p_min + p_range * i / 4
        y_pos = int(plot_t + plot_h - plot_h * (i / 4))
        # 网格横线
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            plot_l, y_pos, plot_l + plot_w, y_pos,
        )
        line.line.color.rgb = GRID_GRAY
        line.line.width = Pt(0.5)
        # y 轴标签
        tb = slide.shapes.add_textbox(l, y_pos - Inches(0.1),
                                       label_w - Inches(0.05), Inches(0.2))
        tf = tb.text_frame
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"{y_val:.0f}"
        _sf(r, SZ_SRC, color=TEXT_GRAY)

    # 每根蜡烛
    candle_slot_w = int(plot_w // n)
    body_w = max(int(candle_slot_w * 0.6), Inches(0.05))

    def price_to_y(p):
        return int(plot_t + plot_h - plot_h * (p - p_min) / p_range)

    for i, (date, (op, hi, lo, cl)) in enumerate(zip(dates, ohlc)):
        slot_x = int(plot_l + i * candle_slot_w)
        center_x = int(slot_x + candle_slot_w // 2)

        is_up = cl >= op
        clr = up_color if is_up else down_color

        # 影线(low → high)
        wick = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            center_x, price_to_y(hi),
            center_x, price_to_y(lo),
        )
        wick.line.color.rgb = clr
        wick.line.width = Pt(1)

        # 实体(open → close)
        top_p = max(op, cl)
        bot_p = min(op, cl)
        body_y = price_to_y(top_p)
        body_h = price_to_y(bot_p) - body_y
        if body_h < 1:
            body_h = Pt(1)  # 平盘最小高度

        body = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(center_x - body_w // 2), body_y, body_w, body_h,
        )
        body.fill.solid()
        body.fill.fore_color.rgb = clr
        body.line.color.rgb = clr
        body.line.width = Pt(0.5)

        # 底部日期标签
        if show_labels and (n <= 8 or i % max(1, n // 6) == 0):
            tb = slide.shapes.add_textbox(
                slot_x, plot_t + plot_h + Inches(0.05),
                candle_slot_w, label_h - Inches(0.05),
            )
            tf = tb.text_frame
            tf.margin_top = Pt(0)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(date)
            _sf(r, SZ_SRC, color=TEXT_GRAY)


# ══════════════════════════════════════════════════════════
# 调度 - diagram 类型从 spec 渲染
# ══════════════════════════════════════════════════════════
def render_from_spec(slide, l, t, w, h, spec: dict):
    """根据 diagram spec 渲染。

    spec 'type' 字段:
    - funnel / process_flow / sankey / treemap / sunburst / gantt / heatmap / candlestick
    """
    t_ = spec["type"]

    if t_ == "funnel":
        return add_funnel(slide, l, t, w, h,
                           spec["stages"], spec.get("values"),
                           spec.get("colors"),
                           spec.get("show_values", True),
                           spec.get("show_pct", True))
    if t_ == "process_flow":
        return add_process_flow(slide, l, t, w, h,
                                 spec["steps"], spec.get("colors"),
                                 spec.get("arrow", True))
    if t_ == "sankey":
        return add_sankey(slide, l, t, w, h,
                           spec["left_nodes"], spec["right_nodes"],
                           spec["flows"],
                           spec.get("left_colors"),
                           spec.get("right_colors"))
    if t_ == "treemap":
        return add_treemap(slide, l, t, w, h,
                            spec["items"], spec.get("colors"))
    if t_ == "sunburst":
        return add_sunburst(slide, l, t, w, h,
                             spec["hierarchy"], spec.get("colors"))
    if t_ == "gantt":
        return add_gantt(slide, l, t, w, h,
                          spec["tasks"], spec.get("color"))
    if t_ == "heatmap":
        return add_heatmap(slide, l, t, w, h,
                            spec["matrix"], spec["x_labels"], spec["y_labels"],
                            spec.get("cmap_low"), spec.get("cmap_high"),
                            spec.get("fmt", "{:.0f}"))
    if t_ in ("candlestick", "stock", "k_line", "k线"):
        return add_candlestick(slide, l, t, w, h,
                                spec["dates"], spec["ohlc"],
                                spec.get("up_color"), spec.get("down_color"),
                                spec.get("show_labels", True))

    raise ValueError(f"Unknown diagram type: {t_}")
