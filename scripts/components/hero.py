"""标题页 / 分节扉页 / 难题扉页 — 大字 Hero 系列。

风格参考:realme 发布会、小米年度战略 deck 的扉页规范。

关键模式:
- add_slide2: 标准内容页(双标题 + 模板占位符复用)
- add_divider: 章节扉页(SECTION 0N + 大字标题 + 灰副标)
- add_problem_divider: 难题扉页(黑底 + 红 badge + 白色大字)
"""

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ..core.font import _sf, FONT_LIGHT, SZ_TITLE, SZ_SUB
from ..core.color import RED, DARK_GRAY, TEXT_GRAY, WHITE
from .text import add_textbox


def add_section_divider(prs, num, title, sub="", layout_idx=2):
    """章节扉页:左对齐红锚 + SECTION 0N 小字 + 大标题 + 灰副标。

    用法:
        add_section_divider(prs, 1, "市场背景", "竞品策略与媒介投放综述")
    """
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # 清空模板占位符
    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)

    # 左侧红锚
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8),
        Inches(0.3), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()

    # SECTION 0N
    add_textbox(slide, Inches(0.8), Inches(2.95), Inches(6), Inches(0.4),
                f"SECTION {num:02d}", sz=Pt(11), bold=True, color=RED,
                align=PP_ALIGN.LEFT, title=True)

    # 大标题
    add_textbox(slide, Inches(0.8), Inches(3.45), Inches(11.5), Inches(1.1),
                title, sz=Pt(36), bold=True, color=DARK_GRAY,
                align=PP_ALIGN.LEFT, title=True)

    # 灰色副标
    if sub:
        add_textbox(slide, Inches(0.8), Inches(4.65), Inches(11.5), Inches(0.6),
                    sub, sz=Pt(14), color=TEXT_GRAY, align=PP_ALIGN.LEFT)
    return slide


def _find_title_placeholder(slide, prefer_idx=None):
    """智能定位 slide 上的标题占位符。

    查找顺序:
    1. 若指定 prefer_idx,优先按 idx 找
    2. type 是 TITLE 或 CENTER_TITLE 的占位符
    3. name 含 "title" / "标题"
    4. idx=0 的占位符(传统约定)
    返回 None 表示找不到。
    """
    from pptx.enum.shapes import PP_PLACEHOLDER

    # 1. 按指定 idx
    if prefer_idx is not None:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == prefer_idx:
                return ph

    # 2. 按 type
    title_types = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
    for ph in slide.placeholders:
        if ph.placeholder_format.type in title_types:
            return ph

    # 3. 按 name
    for ph in slide.placeholders:
        name = (ph.name or "").lower()
        if "title" in name or "标题" in name:
            return ph

    # 4. 按 idx=0
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            return ph

    return None


def add_content_slide(prs, main_title, sub_title="",
                      layout_idx=6, **kwargs):
    """标准内容页 — 用 Blank layout + 自己画标题/副标。

    设计决定:**完全不用 placeholder,所有元素都自己定位**。
    模板的价值在于:
    - 母版字体(已通过 patch_theme_fonts 注入微软雅黑 Light)
    - 色板 / chrome(顶部底部装饰条)
    - 16:9 画布尺寸

    不依赖 placeholder 的好处:
    - 位置完全可控,不会从 master 继承不一致的属性
    - 跨模板兼容性强(只要模板有 layout_idx 这一页就行)
    - LibreOffice / Keynote / PowerPoint 渲染一致

    Args:
        layout_idx: 默认 6 = Blank layout(无任何 placeholder)
        kwargs: 兼容旧调用,所有 placeholder 相关参数被忽略

    Returns:
        (slide, content_top_y) tuple
    """
    from ..core.layout import SAFE_L, FULL_W
    from pptx.enum.text import PP_ALIGN

    # 使用 Blank layout(idx=6 是约定,大多数模板都有)
    # 如果模板的 idx=6 不是 Blank,就降级用最后一个 layout
    try:
        slide_layout = prs.slide_layouts[layout_idx]
    except IndexError:
        slide_layout = prs.slide_layouts[-1]

    slide = prs.slides.add_slide(slide_layout)

    # 删除该 layout 自带的所有占位符(确保画布干净)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    # 主标题 — 固定位置,左对齐
    title_t = Inches(0.35)
    title_h = Inches(0.65)
    add_textbox(slide, SAFE_L, title_t, FULL_W, title_h,
                main_title, sz=SZ_TITLE, color=DARK_GRAY,
                align=PP_ALIGN.LEFT, title=True)

    # 副标 — 紧贴主标下方
    sub_bottom_y = title_t + title_h
    if sub_title:
        sub_t = title_t + title_h + Inches(0.02)
        sub_h = Inches(0.30)
        add_textbox(slide, SAFE_L, sub_t, FULL_W, sub_h,
                    sub_title, sz=SZ_SUB, color=TEXT_GRAY,
                    align=PP_ALIGN.LEFT)
        sub_bottom_y = sub_t + sub_h

    # 标题分隔线(细灰横线)
    from pptx.enum.shapes import MSO_SHAPE
    from ..core.color import RULE_GRAY
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        SAFE_L, sub_bottom_y + Inches(0.05),
        FULL_W, Inches(0.01),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = RULE_GRAY
    rule.line.fill.background()

    # 内容区起始 y(分隔线下方留 0.18in 间距)
    content_top_y = sub_bottom_y + Inches(0.25)
    return slide, content_top_y


def add_problem_divider(prs, num_label, problem_title, big_claim="",
                         layout_idx=2, bg_color=None):
    """难题扉页(黑底 + 红 badge + 白色大字)。

    Args:
        num_label: badge 文字,如 "Q1" / "01"
        problem_title: 主问题(白色大字)
        big_claim: 副标主张(可选)
        bg_color: 自定义背景色,默认黑(#151515)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)

    # 全屏背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color or RGBColor(0x15, 0x15, 0x15)
    bg.line.fill.background()

    # 红 badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.8),
        Inches(1.5), Inches(0.38)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RED
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num_label
    _sf(r, Pt(12), True, WHITE)

    # 大字主张
    add_textbox(slide, Inches(1.0), Inches(3.4), Inches(11), Inches(1.3),
                problem_title, sz=Pt(34), bold=True, color=WHITE, title=True)

    if big_claim:
        add_textbox(slide, Inches(1.0), Inches(4.8), Inches(11), Inches(0.6),
                    big_claim, sz=Pt(14), color=RGBColor(0xCC, 0xCC, 0xCC))
    return slide


def add_hero_title(slide, main_title, sub_title="",
                    main_color=DARK_GRAY, sub_color=TEXT_GRAY):
    """Hero 大字 — 在已有 slide 上添加 32pt 大标题(不创建 slide)。

    适合在自定义 layout 上自由放置 hero 标题。
    """
    from ..core.layout import SAFE_L, FULL_W
    add_textbox(slide, SAFE_L, Inches(2.5), FULL_W, Inches(1.5),
                main_title, sz=Pt(32), bold=True, color=main_color,
                align=PP_ALIGN.LEFT, title=True)
    if sub_title:
        add_textbox(slide, SAFE_L, Inches(4.0), FULL_W, Inches(0.5),
                    sub_title, sz=Pt(14), color=sub_color,
                    align=PP_ALIGN.LEFT)
