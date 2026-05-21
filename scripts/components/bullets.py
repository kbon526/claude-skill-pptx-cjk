"""红点 bullets 组件。

每行 = 独立文本框 + 红色圆点。优于使用段落自带 bullet,因为:
1. 颜色完全可控(段落 bullet 颜色继承 layout,不稳定)
2. 行间距用 y 偏移精确控制
3. 中文场景下圆点和文字基线对齐更可靠
"""

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from ..core.font import _sf, SZ_BODY
from ..core.color import RED, DARK_GRAY


def add_bullets(
    slide, l, t, w, lines,
    spacing=Inches(0.36), sz=SZ_BODY,
    dot_color=RED, text_color=DARK_GRAY,
):
    """红点 bullet 列表。

    Args:
        slide: 目标 slide
        l, t, w: 起始位置和宽度
        lines: list[str] 每行文本
        spacing: 行间距
        sz: 字号
        dot_color: 圆点颜色(默认红)
        text_color: 文本颜色(默认深灰)
    """
    for i, text in enumerate(lines):
        y = t + i * spacing
        # 红圆点(直径 6pt,稍向下偏移对齐文字基线)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, l, y + Pt(5), Pt(6), Pt(6)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
        # 文字框
        tb = slide.shapes.add_textbox(
            l + Pt(14), y, w - Pt(14), Inches(0.32)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        _sf(r, sz, color=text_color)
