"""文本组件:textbox / multiline / source note。

所有文本写入都经过 core.font._sf 注入,保证字体一致。
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from ..core.font import _sf, SZ_BODY, SZ_SRC
from ..core.color import SRC_GRAY
from ..core.layout import SAFE_L, Y_SRC, H_SRC


def add_textbox(
    slide, l, t, w, h, text,
    sz=SZ_BODY, bold=False, color=None,
    align=PP_ALIGN.LEFT, title=False,
):
    """单段文本框。bold/title 是历史兼容参数,实际不使用 Bold 权重。"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if p.runs:
            _sf(p.runs[0], sz, bold, color, title)
    return tb


def add_multiline(slide, l, t, w, h, lines, spacing=Pt(3)):
    """多行混合格式。

    Args:
        lines: list of (text, sz, bold, color, align)
    """
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, sz, bold, clr, al) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        p.space_after = spacing
        r = p.add_run()
        r.text = text
        _sf(r, sz, bold, clr)
    return tb


def add_source_note(slide, text="数据来源: ...", w=Inches(8.8)):
    """页面底部统一来源注。"""
    return add_textbox(
        slide, SAFE_L, Y_SRC, w, H_SRC,
        text, sz=SZ_SRC, color=SRC_GRAY,
    )
