"""KPI 卡片组件。

视觉:浅灰底色矩形 + 大字数值(红) + 小字标签(深灰) + 底部红色细横条。
适合右栏 2×2 网格布局展示关键指标。
"""

from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ..core.font import _sf
from ..core.color import RED, DARK_GRAY, LIGHT_GRAY
from ..core.layout import R_KPI_W, R_KPI_H


def add_kpi(
    slide, l, t,
    w=R_KPI_W, h=R_KPI_H,
    val="", label="",
    value_color=RED, label_color=DARK_GRAY,
    bg_color=LIGHT_GRAY, accent_color=None,
):
    """添加单个 KPI 卡片。

    Args:
        l, t: 左上角位置
        w, h: 宽高(默认用 layout.R_KPI_W/H)
        val: 大字数值,如 "37%" "5.2x"
        label: 小字标签,如 "Toyota Share of Voice"
        value_color: 数值颜色(默认红)
        label_color: 标签颜色(默认深灰)
        bg_color: 卡片底色(默认浅灰)
        accent_color: 底部细横条颜色(默认 = value_color)
    """
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # 内容
    tf = bg.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(val)
    _sf(r, Pt(20), True, value_color, title=True)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(1)
    r2 = p2.add_run()
    r2.text = label
    _sf(r2, Pt(8), color=label_color)

    # 底部细横条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, l, t + h, w, Inches(0.04)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color or value_color
    bar.line.fill.background()
    return bg


def add_kpi_grid_2x2(slide, kpis, top_y, gap_y=Inches(0.85), gap_x=Inches(0.1)):
    """便利方法:右栏 2×2 KPI 网格批量绘制。

    Args:
        kpis: list[dict] 每个 dict 含 val + label + 可选样式参数
    """
    from ..core.layout import kpi_grid_2x2
    positions = kpi_grid_2x2(top=top_y, gap_y=gap_y, gap_x=gap_x)
    shapes = []
    for pos, kpi in zip(positions, kpis):
        l, t, w, h = pos
        shapes.append(add_kpi(slide, l, t, w, h, **kpi))
    return shapes
