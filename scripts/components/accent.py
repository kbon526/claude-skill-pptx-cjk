"""装饰元素:红锚、红 label、列分隔线、insight 卡、recap 块、bento、phase 条、横条。

这些都是中文商务提案常见的小型装饰组件,有固定视觉模式,直接复用避免每次重写。
"""

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ..core.font import _sf, SZ_BODY
from ..core.color import (
    RED, WHITE, DARK_GRAY, TEXT_GRAY, LIGHT_GRAY, RULE_GRAY,
)
from ..core.layout import SAFE_L, FULL_W, COL_GAP_X, Y_CONTENT, GUTTER
from .text import add_textbox


def add_red_label(slide, l, t, w, h, text, color=RED):
    """红色矩形标签(用于章节小标题)。"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    tf = s.text_frame
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _sf(r, Pt(10), True, WHITE)
    return s


def add_red_vline(slide, l, t, h=Inches(1.0), color=RED):
    """红色细竖线(章节强调)。"""
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.04), h)
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()
    return b


def add_col_divider(slide, x=COL_GAP_X, t=Y_CONTENT, h=Inches(5.4)):
    """两栏中缝浅灰竖线。"""
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, t, Inches(0.015), h)
    v.fill.solid()
    v.fill.fore_color.rgb = RULE_GRAY
    v.line.fill.background()
    return v


def add_insight_card(slide, l, t, w, h, title, body, title_color=RED):
    """左红条 + 标题 + 正文 — 单条洞察卡。"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, l, t, Inches(0.05), h
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = title_color
    bar.line.fill.background()
    add_textbox(slide, l + Inches(0.14), t, w - Inches(0.14), Inches(0.26),
                title, sz=Pt(11), bold=True, color=title_color, title=True)
    add_textbox(slide, l + Inches(0.14), t + Inches(0.28),
                w - Inches(0.14), h - Inches(0.28),
                body, sz=SZ_BODY, color=DARK_GRAY)


def add_recap_block(slide, x, y, w, h, num_label, title, body, color):
    """Recap 卡(白底有色边框 + 大数字 + 小标题 + 正文 + 底部短色线)。"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = color
    bg.line.width = Pt(0.75)

    add_textbox(slide, x + Inches(0.25), y + Inches(0.15),
                w - Inches(0.5), Inches(0.5),
                num_label, sz=Pt(28), bold=True, color=color, title=True)
    add_textbox(slide, x + Inches(0.25), y + Inches(0.70),
                w - Inches(0.5), Inches(0.3),
                title, sz=Pt(12), bold=True, color=DARK_GRAY, title=True)
    add_textbox(slide, x + Inches(0.25), y + Inches(1.0),
                w - Inches(0.5), h - Inches(1.1),
                body, sz=SZ_BODY, color=TEXT_GRAY)

    # 底部红色短线
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x + Inches(0.25), y + h - Inches(0.04),
        Inches(0.4), Inches(0.04),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_bento_row(slide, y_top, items, total_w=None, left=None,
                   h=None, gap=None):
    """底部 bento 小卡片行(白底有色边框 + 数字 + 简述)。

    Args:
        items: list of (label, value, desc, color),建议 3 项,最多 4 项
    """
    if total_w is None:
        total_w = FULL_W
    if left is None:
        left = SAFE_L
    if h is None:
        h = Inches(0.6)
    if gap is None:
        gap = Inches(0.12)

    n = len(items)
    bw = (total_w - gap * (n - 1)) / n
    for i, (lbl, val, desc, clr) in enumerate(items):
        bx = left + i * (bw + gap)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, y_top, bw, h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = clr
        bg.line.width = Pt(0.75)
        add_textbox(slide, bx + Inches(0.12), y_top + Inches(0.05),
                    bw - Inches(0.24), Inches(0.2),
                    lbl, sz=Pt(9), color=TEXT_GRAY)
        add_textbox(slide, bx + Inches(0.12), y_top + Inches(0.22),
                    Inches(1.5), Inches(0.35),
                    val, sz=Pt(17), color=clr, title=True)
        add_textbox(slide, bx + Inches(1.65), y_top + Inches(0.24),
                    bw - Inches(1.8), Inches(0.35),
                    desc, sz=Pt(9), color=DARK_GRAY)


def add_phase_track(slide, y_top, phases):
    """横向多阶段 Roadmap 条。

    Args:
        phases: list of (label, title, color),每段一个色块
    """
    n = len(phases)
    x = SAFE_L
    gap = GUTTER
    block_w = (FULL_W - gap * (n - 1)) / n
    for i, (label, title, clr) in enumerate(phases):
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y_top, block_w, Inches(0.65)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = clr
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(5)
        tf.margin_left = Pt(14)
        tf.margin_right = Pt(10)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label
        _sf(r, Pt(9), True, WHITE)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = title
        _sf(r2, Pt(12), True, WHITE, True)
        x += block_w + gap


def add_hbar(slide, l, t, w, label, value, max_val, color, sz_label=SZ_BODY):
    """水平条:左标签 + 中进度条 + 右数值。"""
    add_textbox(slide, l, t, Inches(1.4), Inches(0.28),
                label, sz=sz_label, color=DARK_GRAY)
    track = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        l + Inches(1.45), t + Inches(0.05),
        w - Inches(2.2), Inches(0.18),
    )
    track.fill.solid()
    track.fill.fore_color.rgb = LIGHT_GRAY
    track.line.fill.background()
    fill_w = int((w - Inches(2.2)) * min(value / max_val, 1.0))
    if fill_w > 0:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            l + Inches(1.45), t + Inches(0.05),
            fill_w, Inches(0.18),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
    add_textbox(slide, l + w - Inches(0.75), t, Inches(0.75), Inches(0.28),
                f"{value}", sz=sz_label, bold=True, color=color,
                align=PP_ALIGN.RIGHT)
