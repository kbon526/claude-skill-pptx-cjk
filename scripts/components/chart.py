"""图表组件 — B-Data 路径核心(可编辑的原生 Chart)。

设计原则:
1. 全部使用 python-pptx 原生 Chart 对象,客户在 PowerPoint 里可双击编辑数据
2. 统一字体(微软雅黑 Light)、统一字号、统一灰阶网格线
3. 数据标签默认显示,符合中文商务习惯

A 类:原生 XL_CHART_TYPE 直接支持(双击改 Excel)
- column_chart / bar_chart / grouped_col / stacked_bar(已有)
- pie_chart / doughnut_chart(占比)
- line_chart(趋势)
- area_chart / stacked_area / theme_river(累计趋势)
- radar_chart(多维评分)
- scatter_chart / bubble_chart(关系)
- stock_chart(K线)

B 类(形状组合,数据嵌入代码)在 components/diagrams.py:
- funnel / process_flow / sankey / treemap / sunburst / gantt / heatmap

数据格式约定:
- cats(单系列):list[str | num] — x 轴类目
- vals(单系列):list[num]
- series(多系列):dict[name -> (vals, color)] 或 list[(name, vals, color)]
"""

from pptx.util import Pt
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from ..core.font import FONT_LIGHT, SZ_CAP, SZ_SRC
from ..core.color import GRID_GRAY, TEXT_GRAY, DARK_GRAY


# ══════════════════════════════════════════════════════════
# 内部辅助:统一图表样式
# ══════════════════════════════════════════════════════════
def _chart_font(chart):
    """把图表上所有文字字体改为 微软雅黑 Light。"""
    for attr in ("category_axis", "value_axis"):
        try:
            ax = getattr(chart, attr)
        except (ValueError, AttributeError):
            continue
        if hasattr(ax, "tick_labels"):
            ax.tick_labels.font.size = SZ_SRC
            ax.tick_labels.font.name = FONT_LIGHT
            ax.tick_labels.font.bold = False
            ax.tick_labels.font.color.rgb = TEXT_GRAY
    if chart.has_legend:
        chart.legend.font.size = SZ_SRC
        chart.legend.font.name = FONT_LIGHT
        chart.legend.font.bold = False


def _style_gridlines(chart):
    """主网格线浅灰 + 细线。"""
    try:
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = GRID_GRAY
        va.major_gridlines.format.line.width = Pt(0.5)
    except (ValueError, AttributeError):
        pass


def _style_data_labels(plot, fmt="#,##0", bold=False):
    """数据标签:深灰 + 微软雅黑 Light + 自定义数值格式。"""
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = SZ_CAP
    dl.font.name = FONT_LIGHT
    dl.font.bold = bold
    dl.font.color.rgb = DARK_GRAY
    dl.number_format = fmt


def _legend_bottom(chart):
    """统一图例位置 — 底部、不挤压绘图区。"""
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False


# ══════════════════════════════════════════════════════════
# A 类 1: 柱状图 / 条形图(单系列)
# ══════════════════════════════════════════════════════════
def add_column_chart(slide, l, t, w, h, cats, vals, colors=None, fmt="#,##0"):
    """单系列柱状图(可逐 point 上色)。"""
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("", vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, cd)
    c = cf.chart
    c.has_legend = False
    c.style = 2
    pl = c.plots[0]
    pl.gap_width = 80
    _style_data_labels(pl, fmt)
    if colors:
        s = c.series[0]
        for i, clr in enumerate(colors):
            s.points[i].format.fill.solid()
            s.points[i].format.fill.fore_color.rgb = clr
    _style_gridlines(c)
    _chart_font(c)
    return cf


def add_bar_chart(slide, l, t, w, h, cats, vals, colors=None, fmt="#,##0"):
    """单系列水平条形图。"""
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("", vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, l, t, w, h, cd)
    c = cf.chart
    c.has_legend = False
    c.style = 2
    pl = c.plots[0]
    pl.gap_width = 60
    _style_data_labels(pl, fmt)
    if colors:
        s = c.series[0]
        for i, clr in enumerate(colors):
            s.points[i].format.fill.solid()
            s.points[i].format.fill.fore_color.rgb = clr
    _style_gridlines(c)
    _chart_font(c)
    return cf


# ══════════════════════════════════════════════════════════
# A 类 2: 多系列柱状图(分组 / 堆叠)
# ══════════════════════════════════════════════════════════
def add_grouped_col(slide, l, t, w, h, cats, series_dict):
    """分组柱状图。series_dict: {name: (values, color)}"""
    cd = CategoryChartData()
    cd.categories = cats
    for n, (v, _) in series_dict.items():
        cd.add_series(n, v)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, cd)
    c = cf.chart
    c.style = 2
    _legend_bottom(c)
    pl = c.plots[0]
    pl.gap_width = 80
    for i, (n, (v, clr)) in enumerate(series_dict.items()):
        s = c.series[i]
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = clr
    _style_gridlines(c)
    _chart_font(c)
    return cf


def add_stacked_bar(slide, l, t, w, h, cats, series_list, percent=True):
    """100% 堆叠条形图。series_list: [(name, values, color), ...]"""
    cd = CategoryChartData()
    cd.categories = cats
    for n, v, _ in series_list:
        cd.add_series(n, v)
    chart_type = (XL_CHART_TYPE.BAR_STACKED_100 if percent
                  else XL_CHART_TYPE.BAR_STACKED)
    cf = slide.shapes.add_chart(chart_type, l, t, w, h, cd)
    c = cf.chart
    c.style = 2
    _legend_bottom(c)
    for i, (n, v, clr) in enumerate(series_list):
        s = c.series[i]
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = clr
    _chart_font(c)
    return cf


def add_stacked_col(slide, l, t, w, h, cats, series_list, percent=False):
    """堆叠柱状图(可选 100%)。"""
    cd = CategoryChartData()
    cd.categories = cats
    for n, v, _ in series_list:
        cd.add_series(n, v)
    chart_type = (XL_CHART_TYPE.COLUMN_STACKED_100 if percent
                  else XL_CHART_TYPE.COLUMN_STACKED)
    cf = slide.shapes.add_chart(chart_type, l, t, w, h, cd)
    c = cf.chart
    c.style = 2
    _legend_bottom(c)
    for i, (n, v, clr) in enumerate(series_list):
        s = c.series[i]
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = clr
    _style_gridlines(c)
    _chart_font(c)
    return cf


# ══════════════════════════════════════════════════════════
# A 类 3: 占比图(饼图 / 环形图)
# ══════════════════════════════════════════════════════════
def add_pie_chart(slide, l, t, w, h, cats, vals, colors=None, fmt='0"%"'):
    """饼图。"""
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("", vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.PIE, l, t, w, h, cd)
    c = cf.chart
    c.style = 2
    chart_legend_right(c)
    pl = c.plots[0]
    _style_data_labels(pl, fmt)
    if colors:
        s = c.series[0]
        for i, clr in enumerate(colors):
            s.points[i].format.fill.solid()
            s.points[i].format.fill.fore_color.rgb = clr
    _chart_font(c)
    return cf


def add_doughnut_chart(slide, l, t, w, h, cats, vals, colors, fmt='0"%"'):
    """环形图(占比展示,中心可放主指标)。"""
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("", vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, l, t, w, h, cd)
    c = cf.chart
    c.style = 2
    chart_legend_right(c)
    pl = c.plots[0]
    _style_data_labels(pl, fmt)
    s = c.series[0]
    for i, clr in enumerate(colors):
        s.points[i].format.fill.solid()
        s.points[i].format.fill.fore_color.rgb = clr
    _chart_font(c)
    return cf


def chart_legend_right(chart):
    """图例放右边(饼图/环形图标准位置)。"""
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False


# ══════════════════════════════════════════════════════════
# A 类 4: 趋势图(折线 / 面积)
# ══════════════════════════════════════════════════════════
def add_line_chart(slide, l, t, w, h, cats, series_dict,
                    fmt="#,##0", smooth=False, markers=True):
    """折线图(单/多系列)。

    Args:
        series_dict: {name: (values, color)} — 多系列共用 cats
        smooth: 是否平滑曲线(False=折线,True=曲线)
        markers: 是否显示数据点标记
    """
    cd = CategoryChartData()
    cd.categories = cats
    for n, (v, _) in series_dict.items():
        cd.add_series(n, v)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, l, t, w, h, cd)
    c = cf.chart
    c.style = 2
    if len(series_dict) > 1:
        _legend_bottom(c)
    else:
        c.has_legend = False
    for i, (n, (v, clr)) in enumerate(series_dict.items()):
        s = c.series[i]
        s.format.line.color.rgb = clr
        s.format.line.width = Pt(2.0)
        s.smooth = smooth
        if markers:
            s.marker.style = 8
            s.marker.size = 6
            s.marker.format.fill.solid()
            s.marker.format.fill.fore_color.rgb = clr
    _style_gridlines(c)
    _chart_font(c)
    return cf


def add_area_chart(slide, l, t, w, h, cats, series_dict,
                    stacked=False, percent=False):
    """面积图(单/多系列)。

    Args:
        stacked: 是否堆叠(多系列时建议 True)
        percent: 是否 100% 堆叠(stacked 自动设为 True)
    """
    cd = CategoryChartData()
    cd.categories = cats
    for n, (v, _) in series_dict.items():
        cd.add_series(n, v)
    if percent:
        chart_type = XL_CHART_TYPE.AREA_STACKED_100
    elif stacked:
        chart_type = XL_CHART_TYPE.AREA_STACKED
    else:
        chart_type = XL_CHART_TYPE.AREA
    cf = slide.shapes.add_chart(chart_type, l, t, w, h, cd)
    c = cf.chart
    c.style = 2
    if len(series_dict) > 1:
        _legend_bottom(c)
    else:
        c.has_legend = False
    for i, (n, (v, clr)) in enumerate(series_dict.items()):
        s = c.series[i]
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = clr
    _style_gridlines(c)
    _chart_font(c)
    return cf


def add_theme_river(slide, l, t, w, h, cats, series_dict):
    """主题河流图 — 用 AREA_STACKED 模拟。

    与堆叠面积图视觉差异:中心对称需要 PowerPoint 后期手动调整,
    本函数提供等效的 AREA_STACKED 实现(数据可双击改)。

    若要严格中心对称的 ThemeRiver 效果,需在 PowerPoint 里:
    把 value_axis 改成 "categories in reverse order" 即可。
    """
    return add_area_chart(slide, l, t, w, h, cats, series_dict,
                          stacked=True, percent=False)


# ══════════════════════════════════════════════════════════
# A 类 5: 雷达图(多维评分)
# ══════════════════════════════════════════════════════════
def add_radar_chart(slide, l, t, w, h, cats, series_dict, filled=False):
    """雷达图(多系列)。

    Args:
        series_dict: {name: (values_list, color)}
        filled: 是否填充(False=线状,True=面状)
    """
    cd = CategoryChartData()
    cd.categories = cats
    for n, (v, _) in series_dict.items():
        cd.add_series(n, v)
    chart_type = (XL_CHART_TYPE.RADAR_FILLED if filled
                  else XL_CHART_TYPE.RADAR_MARKERS)
    cf = slide.shapes.add_chart(chart_type, l, t, w, h, cd)
    c = cf.chart
    c.style = 2
    _legend_bottom(c)
    for i, (n, (v, clr)) in enumerate(series_dict.items()):
        s = c.series[i]
        s.format.line.color.rgb = clr
        s.format.line.width = Pt(2.5)
        if not filled:
            s.marker.style = 8
            s.marker.size = 7
            s.marker.format.fill.solid()
            s.marker.format.fill.fore_color.rgb = clr
        else:
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = clr
    _chart_font(c)
    return cf


# ══════════════════════════════════════════════════════════
# A 类 6: 关系图(散点 / 气泡)
# ══════════════════════════════════════════════════════════
def add_scatter_chart(slide, l, t, w, h, series_data, colors=None):
    """散点图(XY)。

    Args:
        series_data: {name: [(x, y), (x, y), ...]} 多组散点
        colors: list[RGBColor] 与 series_data 顺序对应
    """
    cd = XyChartData()
    for name, pts in series_data.items():
        s = cd.add_series(name)
        for x, y in pts:
            s.add_data_point(x, y)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, l, t, w, h, cd)
    c = cf.chart
    c.style = 2
    if len(series_data) > 1:
        _legend_bottom(c)
    else:
        c.has_legend = False
    for i, name in enumerate(series_data.keys()):
        s = c.series[i]
        s.marker.style = 8
        s.marker.size = 7
        if colors and i < len(colors):
            s.marker.format.fill.solid()
            s.marker.format.fill.fore_color.rgb = colors[i]
            s.marker.format.line.color.rgb = colors[i]
    _style_gridlines(c)
    _chart_font(c)
    return cf


def add_bubble_chart(slide, l, t, w, h, series_data, colors=None,
                      bubble_scale=100):
    """气泡图(三维:x、y、size)。

    Args:
        series_data: {name: [(x, y, size), ...]}
        bubble_scale: 气泡大小百分比(默认 100)
    """
    cd = BubbleChartData()
    for name, pts in series_data.items():
        s = cd.add_series(name)
        for x, y, size in pts:
            s.add_data_point(x, y, size)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.BUBBLE, l, t, w, h, cd)
    c = cf.chart
    c.style = 2
    if len(series_data) > 1:
        _legend_bottom(c)
    else:
        c.has_legend = False
    for i, name in enumerate(series_data.keys()):
        s = c.series[i]
        if colors and i < len(colors):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = colors[i]
            s.format.line.color.rgb = colors[i]
    _style_gridlines(c)
    _chart_font(c)
    return cf


# ══════════════════════════════════════════════════════════
# A 类 7: K 线图(股票) — ⚠️ python-pptx 1.x 不支持
# ══════════════════════════════════════════════════════════
# python-pptx 的 STOCK_OHLC 类型 XML 写入未实现,无法用原生 Chart。
# K 线图改为 components/diagrams.py 的 add_candlestick(用形状画)。
# 此处保留接口提示用户。
def add_stock_chart(*args, **kwargs):
    raise NotImplementedError(
        "python-pptx 1.x 的 STOCK_OHLC 写入未实现。\n"
        "请用 components.diagrams.add_candlestick() 代替(基于形状绘制)。"
    )


# ══════════════════════════════════════════════════════════
# 调度 - 根据 spec 自动选图表类型(仅 A 类)
# ══════════════════════════════════════════════════════════
def render_from_spec(slide, l, t, w, h, spec: dict):
    """根据 chart spec 自动渲染原生图表。

    spec 'type' 字段对应:
    - column / bar / grouped_col
    - stacked_col / stacked_bar(可选 percent)
    - pie / doughnut
    - line / area / theme_river
    - radar
    - scatter / bubble
    - stock

    若 type 是 diagram 类(funnel/flow/sankey/...),应该走 diagrams.render_from_spec。
    """
    t_ = spec["type"]

    # 占比类
    if t_ == "pie":
        return add_pie_chart(slide, l, t, w, h,
                             spec["cats"], spec["vals"],
                             spec.get("colors"), spec.get("fmt", '0"%"'))
    if t_ == "doughnut":
        return add_doughnut_chart(slide, l, t, w, h,
                                   spec["cats"], spec["vals"], spec["colors"],
                                   spec.get("fmt", '0"%"'))

    # 比较类
    if t_ == "column":
        return add_column_chart(slide, l, t, w, h,
                                 spec["cats"], spec["vals"],
                                 spec.get("colors"), spec.get("fmt", "#,##0"))
    if t_ == "bar":
        return add_bar_chart(slide, l, t, w, h,
                             spec["cats"], spec["vals"],
                             spec.get("colors"), spec.get("fmt", "#,##0"))
    if t_ == "grouped_col":
        return add_grouped_col(slide, l, t, w, h,
                               spec["cats"], spec["series"])
    if t_ == "stacked_col":
        return add_stacked_col(slide, l, t, w, h,
                                spec["cats"], spec["series_list"],
                                spec.get("percent", False))
    if t_ == "stacked_bar":
        return add_stacked_bar(slide, l, t, w, h,
                                spec["cats"], spec["series_list"],
                                spec.get("percent", True))

    # 趋势类
    if t_ == "line":
        return add_line_chart(slide, l, t, w, h,
                              spec["cats"], spec["series"],
                              spec.get("fmt", "#,##0"),
                              spec.get("smooth", False),
                              spec.get("markers", True))
    if t_ == "area":
        return add_area_chart(slide, l, t, w, h,
                              spec["cats"], spec["series"],
                              spec.get("stacked", False),
                              spec.get("percent", False))
    if t_ == "theme_river":
        return add_theme_river(slide, l, t, w, h,
                                spec["cats"], spec["series"])

    # 多维评分
    if t_ == "radar":
        return add_radar_chart(slide, l, t, w, h,
                                spec["cats"], spec["series"],
                                spec.get("filled", False))

    # 关系类
    if t_ == "scatter":
        return add_scatter_chart(slide, l, t, w, h,
                                  spec["series"], spec.get("colors"))
    if t_ == "bubble":
        return add_bubble_chart(slide, l, t, w, h,
                                 spec["series"], spec.get("colors"),
                                 spec.get("bubble_scale", 100))

    # 金融
    if t_ == "stock":
        raise ValueError(
            "stock 类型走 components.diagrams.render_from_spec(用形状绘制 K 线)"
        )

    raise ValueError(
        f"Unknown chart type: {t_}. "
        f"如果是 diagram 类(funnel/flow/sankey/treemap/sunburst/gantt/heatmap),"
        f"请用 components.diagrams.render_from_spec"
    )
