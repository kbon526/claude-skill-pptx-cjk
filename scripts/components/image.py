"""安全图片插入 — aspect-fit + 居中 + 可选边框 + 可选底部说明。

为什么要"安全":python-pptx 默认 add_picture 会强制拉伸到指定宽高,中文场景
里图片通常需要保持原始比例。本组件:
1. 用 PIL 读取原图尺寸
2. 在指定 bounding box 内做 aspect-fit
3. 居中偏移
4. 文件不存在时降级为带文件名的浅灰占位框
"""

import os
from pathlib import Path

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ..core.font import _sf, SZ_SRC
from ..core.color import LIGHT_GRAY, MED_GRAY, SRC_GRAY
from .text import add_textbox


def add_picture_safe(slide, path, l, t, w, h,
                      caption=None, border=True,
                      placeholder_text=None):
    """安全图片插入。

    Args:
        path: 图片路径
        l, t, w, h: 目标 bounding box
        caption: 底部说明(可选)
        border: 是否加 0.5pt 中灰边框
        placeholder_text: 文件不存在时的占位文字,默认显示 [文件名]
    """
    path_str = str(path)

    if os.path.exists(path_str):
        try:
            from PIL import Image as PILImage
            img = PILImage.open(path_str)
            iw, ih = img.size
            img.close()
        except ImportError:
            # PIL 不可用,降级为强制拉伸
            pic = slide.shapes.add_picture(path_str, l, t, w, h)
            if border:
                pic.line.color.rgb = MED_GRAY
                pic.line.width = Pt(0.5)
            if caption:
                add_textbox(slide, l, t + h + Inches(0.05), w, Inches(0.2),
                            caption, sz=SZ_SRC, color=SRC_GRAY,
                            align=PP_ALIGN.CENTER)
            return pic

        ratio = iw / ih
        box_ratio = w / h
        if box_ratio > ratio:
            actual_h = h
            actual_w = int(h * ratio)
        else:
            actual_w = w
            actual_h = int(w / ratio)

        dx = (w - actual_w) // 2
        dy = (h - actual_h) // 2
        pic = slide.shapes.add_picture(path_str, l + dx, t + dy,
                                       actual_w, actual_h)
        if border:
            pic.line.color.rgb = MED_GRAY
            pic.line.width = Pt(0.5)
        if caption:
            add_textbox(slide, l, t + h + Inches(0.05), w, Inches(0.2),
                        caption, sz=SZ_SRC, color=SRC_GRAY,
                        align=PP_ALIGN.CENTER)
        return pic

    # 文件不存在 — 降级占位框
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.fill.background()
    tf = bg.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = placeholder_text or f"[{os.path.basename(path_str)}]"
    _sf(r, Pt(8), color=MED_GRAY)
    return bg
