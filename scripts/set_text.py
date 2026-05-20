"""
set_text.py — 字体规范封装：统一注入微软雅黑 Bold/Light，避免 shape.text= 直接赋值导致的字体丢失。

设计哲学：
- 标题用 微软雅黑 Bold
- 正文用 微软雅黑 Light（实际字重靠 bold=False 控制）
- 禁用 微软雅黑 Regular
- 中英文用同一字体名（最简方案，避免混排断层）

用法（作为模块）:
    from set_text import set_text, FONT_TITLE, FONT_BODY

    set_text(title_box, "美国市场策略",
             font=FONT_TITLE, size=28, bold=True, color="062032")
    set_text(body_box, "高净值订阅家庭 + 节日采购旺季",
             font=FONT_BODY, size=14, color="333333")

也可以多段落:
    set_paragraphs(text_box, [
        ("第一段标题", {"size": 16, "bold": True}),
        ("第一段内容...", {"size": 12}),
        ("", {}),  # 空行
        ("第二段标题", {"size": 16, "bold": True}),
        ("第二段内容...", {"size": 12}),
    ])
"""

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FONT_TITLE = "微软雅黑"
FONT_BODY = "微软雅黑"

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def set_text(shape, text, *, font=FONT_BODY, size=14,
             bold=False, color="333333", align="left",
             line_spacing=1.2, word_wrap=True):
    """单段落文本写入。"""
    tf = shape.text_frame
    tf.word_wrap = word_wrap
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = ALIGN_MAP.get(align, PP_ALIGN.LEFT)
    p.line_spacing = line_spacing

    run = p.add_run()
    run.text = text
    _apply_font(run, font, size, bold, color)


def set_paragraphs(shape, paragraphs, *, font_default=FONT_BODY,
                   size_default=14, color_default="333333",
                   line_spacing=1.2, word_wrap=True):
    """
    多段落写入。
    paragraphs: List[Tuple[str, dict]]，dict 可包含 size/bold/color/font/align
    """
    tf = shape.text_frame
    tf.word_wrap = word_wrap
    tf.clear()

    for i, (text, opts) in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = ALIGN_MAP.get(opts.get("align", "left"), PP_ALIGN.LEFT)
        p.line_spacing = opts.get("line_spacing", line_spacing)

        if not text:
            continue  # 空行

        run = p.add_run()
        run.text = text
        _apply_font(run,
                    opts.get("font", font_default),
                    opts.get("size", size_default),
                    opts.get("bold", False),
                    opts.get("color", color_default))


def _apply_font(run, font, size, bold, color):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor.from_string(color.lstrip("#"))


def set_title(slide, text, *, size=28, color="062032"):
    """
    写入 slide 的 title 占位符（按 type 而非 idx 找）。

    返回 True 表示成功，False 表示该 slide 无 title 占位符。
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.type in (13, 14):  # TITLE / CENTER_TITLE
            set_text(ph, text, font=FONT_TITLE, size=size,
                     bold=True, color=color, align="left")
            return True
    return False
