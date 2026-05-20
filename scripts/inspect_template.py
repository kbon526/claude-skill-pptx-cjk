"""
inspect_template.py — 解包 PPTX 并枚举 slide_layouts 与占位符索引。

用法：
    python3 inspect_template.py path/to/template.pptx
    python3 inspect_template.py path/to/template.pptx --unpack ./unpacked

输出：
    - 每个 slide_layout 的索引、名称、占位符列表（idx, type, name）
    - 主题字体（major / minor）
    - 母版的装饰元素概览
    - 可选：解包到指定目录

依赖：python-pptx
"""

import argparse
import sys
import zipfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.text import PP_PLACEHOLDER
except ImportError:
    sys.stderr.write(
        "错误：未安装 python-pptx。请运行：pip install python-pptx\n"
        "macOS 用户注意 python3 路径，详见 docs/03-python-pptx-gotchas.md\n"
    )
    sys.exit(1)


PLACEHOLDER_NAMES = {
    PP_PLACEHOLDER.TITLE: "TITLE",
    PP_PLACEHOLDER.BODY: "BODY",
    PP_PLACEHOLDER.CENTER_TITLE: "CENTER_TITLE",
    PP_PLACEHOLDER.SUBTITLE: "SUBTITLE",
    PP_PLACEHOLDER.OBJECT: "OBJECT",
    PP_PLACEHOLDER.PICTURE: "PICTURE",
    PP_PLACEHOLDER.CHART: "CHART",
    PP_PLACEHOLDER.TABLE: "TABLE",
    PP_PLACEHOLDER.MEDIA_CLIP: "MEDIA_CLIP",
    PP_PLACEHOLDER.SLIDE_NUMBER: "SLIDE_NUMBER",
    PP_PLACEHOLDER.DATE: "DATE",
    PP_PLACEHOLDER.FOOTER: "FOOTER",
    PP_PLACEHOLDER.HEADER: "HEADER",
}


def inspect_layouts(prs):
    print("=" * 70)
    print("Slide Layouts")
    print("=" * 70)
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            ph_type = PLACEHOLDER_NAMES.get(ph.placeholder_format.type, "UNKNOWN")
            placeholders.append(
                f"(idx={ph.placeholder_format.idx}, type={ph_type}, name={ph.name!r})"
            )
        ph_str = ", ".join(placeholders) if placeholders else "<空，无占位符>"
        print(f"  Layout[{i}] {layout.name!r}")
        print(f"      {ph_str}")


def inspect_theme_fonts(pptx_path):
    """从 ZIP 包中读 theme1.xml 提取字体定义。"""
    print()
    print("=" * 70)
    print("Theme Fonts")
    print("=" * 70)
    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            theme_xml = z.read("ppt/theme/theme1.xml").decode("utf-8")
    except (KeyError, FileNotFoundError):
        print("  未找到 ppt/theme/theme1.xml")
        return

    from lxml import etree
    root = etree.fromstring(theme_xml.encode("utf-8"))
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    for kind in ("majorFont", "minorFont"):
        font_el = root.find(f".//a:fontScheme/a:{kind}", ns)
        if font_el is None:
            continue
        latin = font_el.find("a:latin", ns)
        ea = font_el.find("a:ea", ns)
        cs = font_el.find("a:cs", ns)
        latin_name = latin.get("typeface") if latin is not None else "?"
        ea_name = ea.get("typeface") if ea is not None else "?"
        cs_name = cs.get("typeface") if cs is not None else "?"
        kind_zh = "标题(major)" if kind == "majorFont" else "正文(minor)"
        print(f"  {kind_zh}: 西文={latin_name!r}, CJK={ea_name!r}, 复杂脚本={cs_name!r}")


def inspect_slide_master(prs):
    print()
    print("=" * 70)
    print("Slide Master Shapes (前 10 个)")
    print("=" * 70)
    master = prs.slide_master
    for i, shape in enumerate(master.shapes):
        if i >= 10:
            print(f"  ... 还有 {len(master.shapes) - 10} 个 shape")
            break
        print(f"  [{i}] {shape.shape_type} {shape.name!r} "
              f"({shape.left}, {shape.top}, {shape.width}, {shape.height})")


def unpack_pptx(pptx_path, dest_dir):
    dest = Path(dest_dir)
    dest.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(pptx_path, 'r') as z:
        z.extractall(dest)
    print(f"\n已解包到 {dest.resolve()}")


def main():
    parser = argparse.ArgumentParser(
        description="枚举 PPTX 模板的 slide_layouts、theme 字体、master 形状",
    )
    parser.add_argument("pptx_path", help="PPTX 文件路径")
    parser.add_argument("--unpack", metavar="DIR",
                        help="额外解包 PPTX 到指定目录")
    args = parser.parse_args()

    pptx_path = Path(args.pptx_path)
    if not pptx_path.is_file():
        sys.stderr.write(f"错误：文件不存在 {pptx_path}\n")
        sys.exit(1)

    prs = Presentation(str(pptx_path))
    print(f"PPTX: {pptx_path.name}")
    print(f"幻灯片数: {len(prs.slides)}")
    print(f"幻灯片尺寸: {prs.slide_width / 914400:.2f} x "
          f"{prs.slide_height / 914400:.2f} 英寸")

    inspect_layouts(prs)
    inspect_theme_fonts(str(pptx_path))
    inspect_slide_master(prs)

    if args.unpack:
        unpack_pptx(str(pptx_path), args.unpack)


if __name__ == "__main__":
    main()
