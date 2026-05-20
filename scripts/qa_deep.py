"""
qa_deep.py — PPTX 结构 QA 脚本。

检查项：
- 形状边界越出 slide（容差 ±0.05 in）
- 文本溢出 textbox（130% 高度容差，含 CJK 宽度估算）
- 标题 / 页脚 placeholder 重叠
- 标题字数过长（>18 字警告）
- run 字号未定义（继承不稳定警告）

用法:
    python3 qa_deep.py output.pptx
    python3 qa_deep.py output.pptx --strict   # 把 warning 也作为 error
    python3 qa_deep.py output.pptx --json     # 输出 JSON 结果

依赖：python-pptx
"""

import argparse
import json
import sys

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.stderr.write("错误：未安装 python-pptx\n")
    sys.exit(1)


EMU_PER_INCH = 914400
TOLERANCE_INCH = 0.05
TEXT_OVERFLOW_RATIO = 1.30
TITLE_MAX_CHARS = 18


def emu_to_inch(emu):
    return emu / EMU_PER_INCH


def is_cjk(c):
    return ('一' <= c <= '鿿'
            or '　' <= c <= '〿'
            or '＀' <= c <= '￯')


def estimate_text_height(text, font_pt, textbox_width_in, line_spacing=1.2):
    """估算文本占用高度（英寸）。"""
    if not text:
        return 0
    cjk = sum(1 for c in text if is_cjk(c))
    latin = len(text) - cjk
    cjk_w = font_pt / 72
    latin_w = font_pt * 0.55 / 72
    total_w = cjk * cjk_w + latin * latin_w
    line_count = max(1, total_w / textbox_width_in) if textbox_width_in else 1
    line_h = font_pt * line_spacing / 72
    return line_count * line_h


def collect_text_runs(text_frame):
    """收集 text_frame 内所有 run 的 text, font_size_pt 列表。"""
    out = []
    for p in text_frame.paragraphs:
        for r in p.runs:
            text = r.text or ""
            sz = r.font.size
            font_pt = sz.pt if sz is not None else None
            out.append((text, font_pt))
    return out


def get_full_text(text_frame):
    return "\n".join(p.text for p in text_frame.paragraphs)


def check_slide(slide, slide_idx, slide_w_in, slide_h_in):
    issues = []  # [(level, message), ...]  level: 'error' | 'warning'

    title_rect = None
    footer_rects = []

    for shape in slide.shapes:
        name = shape.name or "Shape"
        try:
            left = emu_to_inch(shape.left)
            top = emu_to_inch(shape.top)
            width = emu_to_inch(shape.width)
            height = emu_to_inch(shape.height)
        except (TypeError, AttributeError):
            continue

        # 边界检查
        if left < -TOLERANCE_INCH:
            issues.append(("error",
                f'{name} 左边界 {left:.2f}in < 0'))
        if left + width > slide_w_in + TOLERANCE_INCH:
            issues.append(("error",
                f'{name} 右边界 {left + width:.2f}in > {slide_w_in:.2f}in'))
        if top < -TOLERANCE_INCH:
            issues.append(("error",
                f'{name} 上边界 {top:.2f}in < 0'))
        if top + height > slide_h_in + TOLERANCE_INCH:
            issues.append(("error",
                f'{name} 下边界 {top + height:.2f}in > {slide_h_in:.2f}in'))

        # 文本相关检查
        if shape.has_text_frame and shape.text_frame.paragraphs:
            tf = shape.text_frame
            full_text = get_full_text(tf)
            runs = collect_text_runs(tf)

            # run 字号未定义
            for run_text, run_size in runs:
                if run_text.strip() and run_size is None:
                    issues.append(("warning",
                        f'{name} run "{run_text[:20]}…" 字号未定义（依赖继承）'))
                    break  # 一个 shape 报一次就够

            # 估算文本溢出
            if runs and width > 0 and height > 0:
                # 取最大 font_size 作为估算依据
                max_size = max((s for _, s in runs if s), default=14)
                est_h = estimate_text_height(full_text, max_size, width)
                if est_h > height * TEXT_OVERFLOW_RATIO:
                    issues.append(("warning",
                        f'{name} 文本预估高度 {est_h:.2f}in > '
                        f'textbox {height:.2f}in '
                        f'(溢出 {est_h / height * 100:.0f}%)'))

            # title / footer 检测
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (13, 14):  # TITLE / CENTER_TITLE
                    title_rect = (left, top, width, height)
                    if len(full_text) > TITLE_MAX_CHARS:
                        issues.append(("warning",
                            f'标题 "{full_text}" ({len(full_text)} 字) '
                            f'超过 {TITLE_MAX_CHARS} 字'))
                elif ph_type == 15:  # FOOTER
                    footer_rects.append((left, top, width, height))

    # title / footer 重叠检测
    if title_rect:
        for fr in footer_rects:
            if rect_overlap(title_rect, fr):
                issues.append(("error",
                    f'title 与 footer 在 ({fr[0]:.2f}, {fr[1]:.2f}) 重叠'))

    return issues


def rect_overlap(r1, r2):
    l1, t1, w1, h1 = r1
    l2, t2, w2, h2 = r2
    if l1 + w1 <= l2 or l2 + w2 <= l1:
        return False
    if t1 + h1 <= t2 or t2 + h2 <= t1:
        return False
    return True


def get_slide_title(slide):
    for shape in slide.shapes:
        if (shape.is_placeholder
                and shape.placeholder_format.type in (13, 14)
                and shape.has_text_frame):
            return shape.text_frame.text.strip()[:30]
    # fallback：找第一个大字号 textbox
    for shape in slide.shapes:
        if shape.has_text_frame:
            for run_text, run_size in collect_text_runs(shape.text_frame):
                if run_size and run_size >= 24 and len(run_text) < 60:
                    return run_text.strip()[:30]
    return "(无标题)"


def main():
    parser = argparse.ArgumentParser(description="PPTX 结构 QA")
    parser.add_argument("pptx_path")
    parser.add_argument("--strict", action="store_true",
                        help="warning 也算失败（exit 1）")
    parser.add_argument("--json", action="store_true",
                        help="JSON 输出")
    args = parser.parse_args()

    prs = Presentation(args.pptx_path)
    slide_w_in = emu_to_inch(prs.slide_width)
    slide_h_in = emu_to_inch(prs.slide_height)

    total_errors = 0
    total_warnings = 0
    all_results = []

    for i, slide in enumerate(prs.slides):
        title = get_slide_title(slide)
        issues = check_slide(slide, i, slide_w_in, slide_h_in)
        errors = [m for lvl, m in issues if lvl == "error"]
        warnings = [m for lvl, m in issues if lvl == "warning"]
        total_errors += len(errors)
        total_warnings += len(warnings)
        all_results.append({
            "slide": i + 1,
            "title": title,
            "errors": errors,
            "warnings": warnings,
        })

    if args.json:
        print(json.dumps({
            "total_slides": len(prs.slides),
            "errors": total_errors,
            "warnings": total_warnings,
            "results": all_results,
        }, ensure_ascii=False, indent=2))
    else:
        for r in all_results:
            if not r["errors"] and not r["warnings"]:
                continue
            print(f"=== Slide {r['slide']:02d} {r['title']} ===")
            for msg in r["errors"]:
                print(f"  ✗ {msg}")
            for msg in r["warnings"]:
                print(f"  ⚠ {msg}")
            print()
        clean = len(prs.slides) - sum(
            1 for r in all_results if r["errors"] or r["warnings"]
        )
        print(f"Total: {total_errors} errors, "
              f"{total_warnings} warnings, "
              f"{clean}/{len(prs.slides)} slides clean")

    fail_count = total_errors + (total_warnings if args.strict else 0)
    sys.exit(1 if fail_count else 0)


if __name__ == "__main__":
    main()
