"""
render_layout.py — 用 PIL 绘制 PPTX 每张 slide 的 shape 几何边界图。

用途：在 LibreOffice 不可用 / 沙箱受限时做视觉 QA 备用方案。
- 绿框 = placeholder（关键布局锚点）
- 灰框 = 自由形状
- 红字显示 textbox 内文字（前 90 字符）
- 底部蓝色字标注 slide 编号 + layout 名

用法:
    python3 render_layout.py output.pptx qa_layout/
    python3 render_layout.py output.pptx qa_layout/ --dpi 100

依赖：python-pptx, Pillow
"""

import argparse
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    sys.stderr.write("错误：未安装 python-pptx\n")
    sys.exit(1)

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    sys.stderr.write("错误：未安装 Pillow。请运行：pip install Pillow\n")
    sys.exit(1)


EMU_PER_INCH = 914400
DEFAULT_DPI = 80


def emu_to_px(emu, dpi):
    return int(emu / EMU_PER_INCH * dpi)


def get_font(size, prefer_zh=True):
    """尝试加载支持中文的字体，失败回退到默认。"""
    candidates = []
    if prefer_zh:
        candidates += [
            "/System/Library/Fonts/STHeiti Medium.ttc",
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/Hiragino Sans GB.ttc",
            "C:/Windows/Fonts/msyh.ttc",  # 微软雅黑
            "C:/Windows/Fonts/simhei.ttf",
            "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
        ]
    candidates += [
        "/System/Library/Fonts/Helvetica.ttc",
        "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for p in candidates:
        if os.path.exists(p):
            try:
                return ImageFont.truetype(p, size)
            except Exception:
                continue
    return ImageFont.load_default()


def render_slide(slide, slide_idx, slide_w_emu, slide_h_emu,
                 layout_name, output_path, dpi=DEFAULT_DPI):
    w_px = emu_to_px(slide_w_emu, dpi)
    h_px = emu_to_px(slide_h_emu, dpi)
    img = Image.new("RGB", (w_px, h_px + 30), "white")
    draw = ImageDraw.Draw(img)

    label_font = get_font(11)
    text_font = get_font(10)
    footer_font = get_font(11, prefer_zh=False)

    for shape in slide.shapes:
        try:
            left = emu_to_px(shape.left, dpi)
            top = emu_to_px(shape.top, dpi)
            width = emu_to_px(shape.width, dpi)
            height = emu_to_px(shape.height, dpi)
        except (TypeError, AttributeError):
            continue

        is_ph = shape.is_placeholder
        color = (16, 176, 64) if is_ph else (130, 130, 130)
        draw.rectangle(
            [left, top, left + width, top + height],
            outline=color,
            width=2 if is_ph else 1,
        )

        # placeholder 标签
        if is_ph:
            try:
                ph_type = str(shape.placeholder_format.type).split(".")[-1]
                draw.text((left + 3, top + 3), ph_type,
                          fill=(16, 176, 64), font=label_font)
            except AttributeError:
                pass

        # 文字内容
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()[:90]
            if text:
                draw.text((left + 4, top + 18), text,
                          fill=(180, 0, 0), font=text_font)

    # 底部标签
    footer = f"Slide {slide_idx + 1:02d}  |  Layout: {layout_name}"
    draw.rectangle([0, h_px, w_px, h_px + 30], fill=(240, 246, 255))
    draw.text((10, h_px + 7), footer,
              fill=(20, 60, 140), font=footer_font)

    img.save(output_path, "PNG")


def main():
    parser = argparse.ArgumentParser(
        description="PIL 渲染 PPTX 几何布局，用于沙箱环境视觉 QA"
    )
    parser.add_argument("pptx_path")
    parser.add_argument("output_dir")
    parser.add_argument("--dpi", type=int, default=DEFAULT_DPI,
                        help=f"渲染分辨率 px/inch（默认 {DEFAULT_DPI}）")
    args = parser.parse_args()

    out_dir = Path(args.output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(args.pptx_path)
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        out_path = out_dir / f"layout-{i + 1:02d}.png"
        render_slide(slide, i, prs.slide_width, prs.slide_height,
                     layout_name, str(out_path), args.dpi)
        print(f"  ✓ {out_path}")

    print(f"\n共渲染 {len(prs.slides)} 张 slide → {out_dir.resolve()}/")


if __name__ == "__main__":
    main()
