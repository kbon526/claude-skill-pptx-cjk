"""B1 视觉提取 — 从素材里抽取已有图表/图片。

场景:客户给的素材是 PDF/PPTX,里面有现成的图表、表格、流程图。
我们想"原样保留"这些视觉,而不是重画。

实现策略:
- PDF: 用 pdfplumber 提取每页的图片,导出为 PNG
- PPTX: 遍历 shapes,把 picture / chart / table 截图导出
- 输出:work/extracted/<source>/<slide_X>_<shape_Y>.png

注:PPTX 的 chart 难以直接"截图",通常的做法是把整页转图后裁剪。
本模块只做"图片"提取,Chart 提取由 chart_remake 重画路径处理。
"""

from pathlib import Path
from typing import Iterable


def extract_pdf_images(pdf_path: str | Path, out_dir: str | Path) -> list[Path]:
    """从 PDF 提取所有嵌入图片。需要 pdfplumber。"""
    try:
        import pdfplumber
    except ImportError:
        raise RuntimeError("Need `pip install pdfplumber Pillow`")

    pdf_path = Path(pdf_path)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    extracted: list[Path] = []
    with pdfplumber.open(str(pdf_path)) as pdf:
        for page_idx, page in enumerate(pdf.pages, 1):
            for img_idx, img in enumerate(page.images, 1):
                # pdfplumber 提供原始坐标 + bytestream
                # 对于嵌入图片,我们用 page.crop + to_image 截取
                bbox = (img["x0"], img["top"], img["x1"], img["bottom"])
                cropped = page.crop(bbox).to_image(resolution=200)
                out_path = out_dir / f"pdf_page{page_idx:02d}_img{img_idx:02d}.png"
                cropped.save(str(out_path))
                extracted.append(out_path)
    return extracted


def extract_pptx_images(pptx_path: str | Path, out_dir: str | Path) -> list[Path]:
    """从 PPTX 提取所有 picture shape。需要 python-pptx。"""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("Need `pip install python-pptx`")

    pptx_path = Path(pptx_path)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    extracted: list[Path] = []
    prs = Presentation(str(pptx_path))
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape_idx, shape in enumerate(slide.shapes, 1):
            if not shape.shape_type:
                continue
            # 13 = MSO_SHAPE_TYPE.PICTURE
            if shape.shape_type == 13:
                ext = shape.image.ext
                out_path = out_dir / f"pptx_slide{slide_idx:02d}_pic{shape_idx:02d}.{ext}"
                with open(out_path, "wb") as f:
                    f.write(shape.image.blob)
                extracted.append(out_path)
    return extracted


def extract_all(sources: Iterable[str | Path], out_dir: str | Path) -> dict[str, list[Path]]:
    """批量提取,返回 {source_name: [extracted_paths]}。"""
    out_dir = Path(out_dir)
    result = {}
    for src in sources:
        src = Path(src)
        sub_out = out_dir / src.stem
        suffix = src.suffix.lower()
        if suffix == ".pdf":
            result[src.name] = extract_pdf_images(src, sub_out)
        elif suffix == ".pptx":
            result[src.name] = extract_pptx_images(src, sub_out)
        else:
            print(f"⚠️ Skip unsupported source: {src}")
            result[src.name] = []
    return result


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 3:
        print("Usage: python visual_extract.py <source.pdf|.pptx> <out_dir>")
        sys.exit(1)
    src = sys.argv[1]
    out = sys.argv[2]
    suffix = Path(src).suffix.lower()
    if suffix == ".pdf":
        files = extract_pdf_images(src, out)
    elif suffix == ".pptx":
        files = extract_pptx_images(src, out)
    else:
        print(f"Unsupported: {suffix}")
        sys.exit(1)
    print(f"Extracted {len(files)} files to {out}/")
    for f in files:
        print(f"  - {f}")
