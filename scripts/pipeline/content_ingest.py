"""A2 内容摄入 — 将各种素材转为统一 Markdown。

支持:
- PDF (借助 markitdown 或 pdfplumber)
- DOCX (借助既有 scripts/extract_docx.py 或 markitdown)
- PPTX (借助 markitdown)
- HTML / 纯文本(直接读)

输出:统一的 Markdown 字符串,带最简单的章节标记(# / ##)。
后续由 content_allocate 切片到 slide。

为什么不直接用 markitdown 库:
1. markitdown 的 PDF 提取对中文有时表现差
2. 我们已有 extract_docx.py 对 DOCX 做了零依赖处理
3. 提供统一接口,后端可热替换
"""

import subprocess
from pathlib import Path
from typing import Iterable


def ingest_one(path: str | Path) -> str:
    """读取单个文件,返回 markdown 字符串。"""
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"Not found: {p}")

    suffix = p.suffix.lower()
    if suffix in (".md", ".markdown", ".txt"):
        return p.read_text(encoding="utf-8")
    if suffix == ".docx":
        return _ingest_docx(p)
    if suffix == ".pdf":
        return _ingest_pdf(p)
    if suffix == ".pptx":
        return _ingest_pptx(p)
    raise ValueError(f"Unsupported format: {suffix}")


def ingest_many(paths: Iterable[str | Path]) -> dict[str, str]:
    """批量摄入,返回 {filename: markdown}。"""
    return {Path(p).name: ingest_one(p) for p in paths}


def _ingest_docx(path: Path) -> str:
    """用既有 scripts/extract_docx.py 解压 DOCX。

    优先调用本仓库的 extract_docx.py(零依赖),
    fallback 到 python-docx 或 markitdown。
    """
    # 尝试本仓库的 extract_docx.py
    extract_script = Path(__file__).parent.parent / "extract_docx.py"
    if extract_script.exists():
        try:
            result = subprocess.run(
                ["python3", str(extract_script), str(path)],
                capture_output=True, text=True, check=True,
            )
            return result.stdout
        except subprocess.CalledProcessError:
            pass

    # Fallback: python-docx
    try:
        from docx import Document
        doc = Document(str(path))
        lines = []
        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                continue
            style = p.style.name if p.style else ""
            if "Heading 1" in style:
                lines.append(f"# {text}")
            elif "Heading 2" in style:
                lines.append(f"## {text}")
            elif "Heading 3" in style:
                lines.append(f"### {text}")
            else:
                lines.append(text)
        return "\n\n".join(lines)
    except ImportError:
        raise RuntimeError(
            "DOCX ingestion needs either scripts/extract_docx.py or `pip install python-docx`"
        )


def _ingest_pdf(path: Path) -> str:
    """PDF 摄入。优先 markitdown,fallback pdfplumber。"""
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        return md.convert(str(path)).text_content
    except ImportError:
        pass
    try:
        import pdfplumber
        chunks = []
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                txt = page.extract_text() or ""
                if txt.strip():
                    chunks.append(f"## Page {i}\n\n{txt}")
        return "\n\n".join(chunks)
    except ImportError:
        raise RuntimeError(
            "PDF ingestion needs `pip install markitdown` or `pip install pdfplumber`"
        )


def _ingest_pptx(path: Path) -> str:
    """PPTX 摄入。优先 markitdown,fallback python-pptx。"""
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        return md.convert(str(path)).text_content
    except ImportError:
        pass
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        chunks = []
        for i, slide in enumerate(prs.slides, 1):
            chunks.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        text = p.text.strip()
                        if text:
                            chunks.append(text)
        return "\n\n".join(chunks)
    except ImportError:
        raise RuntimeError(
            "PPTX ingestion needs `pip install markitdown` or `pip install python-pptx`"
        )


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python content_ingest.py <file>")
        sys.exit(1)
    print(ingest_one(sys.argv[1]))
