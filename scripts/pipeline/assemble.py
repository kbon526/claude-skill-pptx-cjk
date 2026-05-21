"""C3 组装引擎 — Asset Registry → .pptx。

⚠️ 强约束:必须提供真实 PPTX 模板(template_path 必填)。
不再支持空白创建 —— python-pptx 默认空白 Presentation 的 layout / 占位符
约定与中文商务模板差异大,渲染效果不可控。

正确用法:
    assemble(reg, template_path="templates/starter.pptx", output_path="out.pptx")

如果你没有模板:
    python3 scripts/tools/generate_starter_template.py templates/starter.pptx
    # 生成本 skill 兼容的最简模板

如果你想检查既有模板是否兼容:
    python3 scripts/tools/validate_template.py path/to/your.pptx

布局映射(layout 字段 → builder):
- section_divider → 章节扉页(全屏 hero 风格)
- problem_divider → 难题扉页(黑底 + 红 badge)
- two_column_kpi_bullets → 左 bullets + 右 KPI 网格(无图表,避免重叠)
- kpi_grid → 全宽 KPI 网格
- bullets → 全宽 bullets(从标题下方开始)
- paragraph → 全宽段落
- chart_focus → 主标题 + 大图表
- two_col_chart → 左 bullets + 右 chart
- default → 自动判断
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ..core.font import patch_theme_fonts, SZ_BODY
from ..core.registry import all_approved, validate
from ..core.layout import (
    SAFE_L, FULL_W, Y_CONTENT, Y_KPI1,
    LW, RX, RW,
)
from ..components import text as comp_text
from ..components import bullets as comp_bullets
from ..components import kpi as comp_kpi
from ..components import chart as comp_chart
from ..components import diagrams as comp_diagrams
from ..components import hero as comp_hero
from ..components import image as comp_image


# ══════════════════════════════════════════════════════════
# 视觉调度:根据 visual.kind 分流到 chart 或 diagrams
# ══════════════════════════════════════════════════════════
def _render_visual(slide, l, t, w, h, visual: dict):
    """根据 kind 调度到 chart 或 diagrams 模块。"""
    kind = visual.get("kind")
    spec = visual.get("spec")
    if kind == "native_chart":
        return comp_chart.render_from_spec(slide, l, t, w, h, spec)
    if kind == "diagram":
        return comp_diagrams.render_from_spec(slide, l, t, w, h, spec)
    if kind == "raster_png":
        return comp_image.add_picture_safe(
            slide, visual["path"], l, t, w, h,
            caption=visual.get("caption"),
        )
    if kind == "image_extracted":
        return comp_image.add_picture_safe(
            slide, visual["path"], l, t, w, h,
            caption=visual.get("caption"),
        )
    return None


# ══════════════════════════════════════════════════════════
# 模板校验
# ══════════════════════════════════════════════════════════
def validate_template(template_path: str | Path) -> list[str]:
    """检查模板是否含本 skill 期望的 layout 结构。

    Returns:
        warnings: 警告列表(空表示完全兼容)
    """
    warnings = []
    prs = Presentation(str(template_path))

    n_layouts = len(prs.slide_layouts)
    if n_layouts < 2:
        warnings.append(f"模板 layout 数量过少({n_layouts}),建议至少 2 个")

    # 期望:存在至少一个 layout 含标题占位符
    has_title_layout = False
    from pptx.enum.shapes import PP_PLACEHOLDER
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if ph.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
            ):
                has_title_layout = True
                break
        if has_title_layout:
            break
    if not has_title_layout:
        warnings.append("没有任何 layout 含标题占位符,主标题将走 fallback 路径(add_textbox)")

    # 检查画布尺寸
    width_in = prs.slide_width / 914400
    height_in = prs.slide_height / 914400
    ratio = width_in / height_in
    if abs(ratio - 16 / 9) > 0.05:
        warnings.append(
            f"模板画布 {width_in:.2f}x{height_in:.2f}in (比例 {ratio:.2f}),"
            f"本 skill 网格常数按 16:9 设计(13.33x7.5),可能错位"
        )

    return warnings


# ══════════════════════════════════════════════════════════
# 布局 builders
# ══════════════════════════════════════════════════════════
def _build_section_divider(prs, slide_data):
    """章节扉页。

    若 visuals 含 raster_png(AI 生成的 hero 图),用它作为全屏背景,
    标题和副标叠加在图上(白字 + 半透明深色蒙版,保证可读性)。
    否则走默认无图版本(白底 + SECTION 0N + 大标题)。
    """
    title = slide_data.get("title", "")
    sub = slide_data.get("subtitle", "") or slide_data.get("copy", {}).get("sub", "")

    # 检测是否有 hero 背景图
    visuals = slide_data.get("visuals", [])
    bg_visual = next(
        (v for v in visuals if v.get("kind") == "raster_png"), None
    )

    if bg_visual:
        return _build_section_divider_with_hero(prs, title, sub, bg_visual)

    # 默认无图版本(沿用 add_section_divider 的 SECTION 0N 风格)
    return comp_hero.add_section_divider(prs, 1, title, sub)


def _build_section_divider_with_hero(prs, title, sub, bg_visual):
    """带 hero 全屏图的分节扉页。

    布局:
    - 全屏 hero 图作为背景(填满 slide)
    - 半透明深色蒙版(底部覆盖,保证白字可读)
    - 大标题(白色)+ 副标(浅灰)在底部 1/3 处
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    from ..core.color import WHITE, RGBColor
    from ..core.font import _sf
    from ..core.layout import SAFE_L, FULL_W, SLIDE_W, SLIDE_H

    # 用空白 layout
    try:
        slide_layout = prs.slide_layouts[6]  # Blank
    except IndexError:
        slide_layout = prs.slide_layouts[-1]
    slide = prs.slides.add_slide(slide_layout)

    # 删掉所有占位符
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    # 1. 全屏 hero 图(图已经过 PIL 后处理裁切到精确 16:9,所以直接铺满不变形)
    img_path = bg_visual.get("path")
    if img_path and Path(img_path).exists():
        # 探测图片实际比例;若已是 16:9 ± 1%,直接铺满;否则用 add_picture_safe 居中
        try:
            from PIL import Image
            with Image.open(img_path) as im:
                iw, ih = im.size
            ratio = iw / ih
            slide_ratio = prs.slide_width / prs.slide_height  # ≈ 1.778
            if abs(ratio - slide_ratio) < 0.02:
                # 比例匹配,直接铺满(不拉伸,因为本来就是 16:9)
                slide.shapes.add_picture(
                    img_path, 0, 0, prs.slide_width, prs.slide_height
                )
            else:
                # 比例不匹配,用 aspect-fit 居中(避免拉伸失真)
                comp_image.add_picture_safe(
                    slide, img_path,
                    0, 0, prs.slide_width, prs.slide_height,
                    border=False,
                )
        except ImportError:
            slide.shapes.add_picture(
                img_path, 0, 0, prs.slide_width, prs.slide_height
            )

    # 2. 底部半透明深色蒙版(从 60% 高度往下,渐黑)
    mask_y = int(prs.slide_height * 0.55)
    mask_h = prs.slide_height - mask_y
    mask = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, mask_y, prs.slide_width, mask_h
    )
    mask.fill.solid()
    mask.fill.fore_color.rgb = RGBColor(0x10, 0x10, 0x10)
    # python-pptx 不直接支持 alpha,用接近黑的纯色作为半透明替代
    mask.line.fill.background()

    # 3. 大标题(白色,底部居中)
    title_y = int(prs.slide_height * 0.62)
    tb = slide.shapes.add_textbox(
        SAFE_L, title_y, FULL_W, Inches(1.4)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    _sf(r, Pt(36), True, WHITE, title=True)

    # 4. 副标(浅灰)
    if sub:
        sub_y = title_y + Inches(1.0)
        tb2 = slide.shapes.add_textbox(
            SAFE_L, sub_y, FULL_W, Inches(0.5)
        )
        tf2 = tb2.text_frame
        tf2.margin_top = Pt(0)
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = sub
        _sf(r2, Pt(14), False, RGBColor(0xCC, 0xCC, 0xCC))

    return slide


def _build_problem_divider(prs, slide_data):
    title = slide_data.get("title", "")
    claim = slide_data.get("subtitle", "")
    badge = slide_data.get("copy", {}).get("badge", "?")
    return comp_hero.add_problem_divider(prs, badge, title, claim)


def _build_two_column_kpi_bullets(prs, slide_data):
    """两栏:左 bullets + 右 2×2 KPI 网格(不放图表,避免重叠)。"""
    slide, content_y = comp_hero.add_content_slide(
        prs, slide_data.get("title", ""), slide_data.get("subtitle", "")
    )
    copy = slide_data.get("copy", {})
    bullets = copy.get("bullets", [])
    if bullets:
        comp_bullets.add_bullets(
            slide, SAFE_L, content_y, LW, bullets[:6],
        )
    kpis_raw = copy.get("kpis", [])
    if kpis_raw:
        kpi_dicts = [_split_kpi_text(k) for k in kpis_raw[:4]]
        comp_kpi.add_kpi_grid_2x2(slide, kpi_dicts, top_y=content_y)
    # ⚠️ 此 layout 不画 chart 以避免与 KPI 重叠
    # 如果需要图表,改用 layout=two_col_chart
    return slide


def _build_two_col_chart(prs, slide_data):
    """两栏:左 bullets + 右 chart/diagram(图表替代 KPI)。"""
    slide, content_y = comp_hero.add_content_slide(
        prs, slide_data.get("title", ""), slide_data.get("subtitle", "")
    )
    copy = slide_data.get("copy", {})
    bullets = copy.get("bullets", [])
    if bullets:
        comp_bullets.add_bullets(
            slide, SAFE_L, content_y, LW, bullets[:6],
        )
    visuals = slide_data.get("visuals", [])
    visual = next((v for v in visuals
                   if v.get("kind") in ("native_chart", "diagram")), None)
    if visual:
        chart_h = Inches(7.0) - content_y - Inches(0.3)
        _render_visual(slide, RX, content_y, RW, chart_h, visual)
    return slide


def _build_kpi_grid(prs, slide_data):
    """全宽 KPI 网格(最多 4 个)。"""
    slide, content_y = comp_hero.add_content_slide(
        prs, slide_data.get("title", ""), slide_data.get("subtitle", "")
    )
    copy = slide_data.get("copy", {})
    kpis_raw = copy.get("kpis", [])
    kpi_dicts = [_split_kpi_text(k) for k in kpis_raw[:4]]
    if kpi_dicts:
        from ..core.layout import kpi_grid_2x2
        positions = kpi_grid_2x2(top=content_y)
        for pos, kdict in zip(positions, kpi_dicts):
            l, t, w, h = pos
            comp_kpi.add_kpi(slide, l, t, w, h, **kdict)
    return slide


def _build_bullets(prs, slide_data):
    """全宽 bullets(从内容起始 y)。"""
    slide, content_y = comp_hero.add_content_slide(
        prs, slide_data.get("title", ""), slide_data.get("subtitle", "")
    )
    copy = slide_data.get("copy", {})
    bullets = copy.get("bullets", [])
    if bullets:
        comp_bullets.add_bullets(
            slide, SAFE_L, content_y, FULL_W, bullets,
            spacing=Inches(0.5),
        )
    return slide


def _build_paragraph(prs, slide_data):
    """全宽段落。"""
    slide, content_y = comp_hero.add_content_slide(
        prs, slide_data.get("title", ""), slide_data.get("subtitle", "")
    )
    copy = slide_data.get("copy", {})
    main = copy.get("main", "")
    if main:
        comp_text.add_textbox(
            slide, SAFE_L, content_y, FULL_W, Inches(7.0) - content_y - Inches(0.3),
            main, sz=SZ_BODY,
        )
    return slide


def _build_chart_focus(prs, slide_data):
    """主标题 + 大图表/diagram(占下方所有可用空间)。"""
    slide, content_y = comp_hero.add_content_slide(
        prs, slide_data.get("title", ""), slide_data.get("subtitle", "")
    )
    visuals = slide_data.get("visuals", [])
    visual = next((v for v in visuals
                   if v.get("kind") in ("native_chart", "diagram")), None)
    if visual:
        chart_h = Inches(7.0) - content_y - Inches(0.3)
        _render_visual(slide, SAFE_L, content_y, FULL_W, chart_h, visual)
    return slide


def _build_default(prs, slide_data):
    """fallback — 万能内容页(段落 + 视觉)。"""
    slide, content_y = comp_hero.add_content_slide(
        prs, slide_data.get("title", ""), slide_data.get("subtitle", "")
    )
    copy = slide_data.get("copy", {})
    main = copy.get("main", "")
    if main:
        comp_text.add_textbox(
            slide, SAFE_L, content_y, FULL_W, Inches(7.0) - content_y - Inches(0.3),
            main, sz=SZ_BODY,
        )
    return slide


# ══════════════════════════════════════════════════════════
# 辅助
# ══════════════════════════════════════════════════════════
def _split_kpi_text(raw) -> dict:
    """把 '37% Toyota Share' 拆成 val + label,或直接接受 dict。"""
    if isinstance(raw, dict):
        return raw
    parts = str(raw).split(maxsplit=1)
    if len(parts) == 2:
        return {"val": parts[0], "label": parts[1]}
    return {"val": str(raw), "label": ""}


# ══════════════════════════════════════════════════════════
# 调度表
# ══════════════════════════════════════════════════════════
LAYOUT_BUILDERS = {
    "section_divider": _build_section_divider,
    "problem_divider": _build_problem_divider,
    "two_column_kpi_bullets": _build_two_column_kpi_bullets,
    "two_col_chart": _build_two_col_chart,
    "kpi_grid": _build_kpi_grid,
    "bullets": _build_bullets,
    "paragraph": _build_paragraph,
    "chart_focus": _build_chart_focus,
    "default": _build_default,
}


# ══════════════════════════════════════════════════════════
# 主入口
# ══════════════════════════════════════════════════════════
def assemble(reg: dict,
             template_path: str | Path,
             output_path: str | Path,
             skip_checkpoint_check: bool = False,
             skip_template_validate: bool = False) -> Path:
    """组装 Asset Registry → .pptx。

    Args:
        reg: Asset Registry dict
        template_path: ⭐ 必填,真实 PPTX 模板路径
        output_path: 输出 .pptx 路径
        skip_checkpoint_check: 是否跳过 checkpoint 检查
        skip_template_validate: 是否跳过模板兼容性校验

    Raises:
        ValueError: registry 无效 / 模板路径不存在
        RuntimeError: checkpoint 未全部通过
    """
    if not template_path:
        raise ValueError(
            "template_path 是必填参数。本 skill 强制要求基于真实 PPTX 模板。\n"
            "如果你没有现成模板,运行:\n"
            "  python3 scripts/tools/generate_starter_template.py templates/starter.pptx\n"
            "生成本 skill 兼容的最简模板。"
        )
    template_path = Path(template_path)
    if not template_path.exists():
        raise FileNotFoundError(f"模板文件不存在: {template_path}")

    errors = validate(reg)
    if errors:
        raise ValueError("Registry validation failed:\n  " + "\n  ".join(errors))

    if not skip_checkpoint_check and not all_approved(reg):
        statuses = reg["meta"]["checkpoint_status"]
        pending = [k for k, v in statuses.items() if v != "approved"]
        raise RuntimeError(
            f"Cannot assemble: pending checkpoints {pending}.\n"
            "Use skip_checkpoint_check=True only if you really know what you're doing."
        )

    # 模板兼容性校验
    if not skip_template_validate:
        warnings = validate_template(template_path)
        if warnings:
            print("⚠️ 模板兼容性警告:")
            for w in warnings:
                print(f"   - {w}")
            print("   (skill 会尽力 fallback,但建议先运行 validate_template.py 工具)")

    prs = Presentation(str(template_path))

    # 强制中文字体
    patch_theme_fonts(prs)

    # 清空模板既有 slides(保留 layouts)
    sld_id_lst = prs.slides._sldIdLst
    for s_id in list(sld_id_lst):
        rId = s_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rId:
            prs.part.drop_rel(rId)
        sld_id_lst.remove(s_id)

    # 按 slide_id 顺序构建
    for sid in sorted(reg["slides"].keys()):
        slide_data = reg["slides"][sid]
        layout = slide_data.get("layout", "default")
        builder = LAYOUT_BUILDERS.get(layout, _build_default)
        try:
            builder(prs, slide_data)
        except Exception as e:
            print(f"⚠️ Failed to build {sid} (layout={layout}): {e}")

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"✅ Saved: {output_path}")
    return output_path


if __name__ == "__main__":
    import sys
    from ..core.registry import load
    if len(sys.argv) < 4:
        print("Usage: python assemble.py <registry.json> <template.pptx> <output.pptx>")
        sys.exit(1)
    reg = load(sys.argv[1])
    assemble(reg, template_path=sys.argv[2], output_path=sys.argv[3],
             skip_checkpoint_check=True)
