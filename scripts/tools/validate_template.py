"""模板兼容性校验工具。

用法:
    python3 scripts/tools/validate_template.py path/to/your.pptx

输出:
- ✅ 模板含 N 个 layout
- ✅ 标题占位符状态
- ⚠️ 不兼容项的具体提示

调用方式:
    /Library/Frameworks/Python.framework/Versions/3.13/bin/python3 \\
        scripts/tools/validate_template.py templates/starter.pptx
"""

import sys
from pathlib import Path

# 把 scripts 加入路径
THIS_DIR = Path(__file__).parent
ROOT = THIS_DIR.parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


def validate(template_path: str) -> int:
    """返回 0=完全兼容 / 1=有警告 / 2=致命问题。"""
    p = Path(template_path)
    if not p.exists():
        print(f"❌ 文件不存在: {p}")
        return 2

    try:
        prs = Presentation(str(p))
    except Exception as e:
        print(f"❌ 无法加载: {e}")
        return 2

    print(f"📂 模板: {p}")
    print(f"📐 画布: {prs.slide_width/914400:.2f}in × {prs.slide_height/914400:.2f}in")

    # 比例检查
    ratio = prs.slide_width / prs.slide_height
    if abs(ratio - 16/9) < 0.05:
        print(f"   ✅ 16:9 比例(本 skill 标准)")
    elif abs(ratio - 4/3) < 0.05:
        print(f"   ⚠️ 4:3 比例(本 skill 网格按 16:9 设计,可能错位)")
    else:
        print(f"   ⚠️ 非标准比例 {ratio:.2f}(预期 16:9 = 1.78)")

    n_layouts = len(prs.slide_layouts)
    print(f"📑 Layouts: {n_layouts} 个")

    # 列出所有 layout + 标题占位符状态
    has_title = False
    for i, layout in enumerate(prs.slide_layouts):
        title_phs = []
        for ph in layout.placeholders:
            if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                title_phs.append(f"idx={ph.placeholder_format.idx}")
        marker = "✅" if title_phs else "  "
        if title_phs:
            has_title = True
        print(f"   {marker} layout[{i}] {layout.name!r}  "
              f"placeholders={len(list(layout.placeholders))}  "
              f"title={title_phs or '无'}")

    print()

    # 综合判断
    warnings = 0
    if not has_title:
        print("❌ 严重:没有任何 layout 含标题占位符")
        print("   → 主标题会走 fallback 路径(直接 add_textbox)")
        print("   → 建议生成 starter 模板:")
        print("     python3 scripts/tools/generate_starter_template.py templates/starter.pptx")
        warnings += 2

    if abs(ratio - 16/9) > 0.05:
        warnings += 1

    if warnings == 0:
        print("✅ 模板完全兼容,可直接使用")
        return 0
    elif warnings == 1:
        print("⚠️ 模板基本可用,有部分警告(渲染可能略有错位)")
        return 1
    else:
        print("❌ 模板不兼容,建议替换或生成 starter 模板")
        return 2


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    sys.exit(validate(sys.argv[1]))
