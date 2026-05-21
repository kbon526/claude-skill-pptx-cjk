"""字体注入器冒烟测试。"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt

from scripts.core.font import (
    _sf, patch_theme_fonts, write_run,
    FONT_LIGHT, SZ_TITLE, SZ_BODY, SZ_HERO,
)


def test_font_constants_exist():
    """字体常量必须可访问。"""
    assert FONT_LIGHT == "微软雅黑 Light"
    assert SZ_TITLE.pt == 22
    assert SZ_BODY.pt == 10
    assert SZ_HERO.pt == 32


def test_sf_injects_font_to_run():
    """_sf 应该正确设置 run 的字体属性。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "测试中文字体"
    _sf(r, SZ_BODY)

    assert r.font.name == FONT_LIGHT
    assert r.font.size.pt == 10
    # bold 必须为 False(本 skill 规则)
    assert r.font.bold is False


def test_patch_theme_fonts_changes_theme():
    """patch_theme_fonts 应该修改主题字体到微软雅黑 Light。"""
    prs = Presentation()
    patch_theme_fonts(prs)

    # 检查主题 part 内容
    for part in prs.part.package.iter_parts():
        if "theme" in part.partname.lower():
            blob = part.blob.decode("utf-8", errors="ignore")
            assert "微软雅黑 Light" in blob


def test_write_run_helper():
    """write_run 应当能写入并注入字体。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    p = tb.text_frame.paragraphs[0]

    r = write_run(p, "Hello 中文", size=SZ_TITLE)
    assert r.text == "Hello 中文"
    assert r.font.name == FONT_LIGHT
    assert r.font.size.pt == 22


if __name__ == "__main__":
    # 不依赖 pytest,直接跑
    for name, fn in list(globals().items()):
        if name.startswith("test_") and callable(fn):
            try:
                fn()
                print(f"✅ {name}")
            except Exception as e:
                print(f"❌ {name}: {e}")
