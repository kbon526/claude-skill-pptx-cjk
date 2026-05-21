"""18 种图表类型冒烟测试 — chart.py + diagrams.py + chart_picker.py。"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.util import Inches

from scripts.components import chart, diagrams
from scripts.core.color import RED, BLUE, GREEN, ORANGE, BLACK
from scripts.pipeline import chart_picker


def make_blank():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


# ══════════════════════════════════════════════════════════
# A 类:原生 Chart 测试(共 12 个)
# ══════════════════════════════════════════════════════════
def test_chart_column():
    _, s = make_blank()
    cf = chart.add_column_chart(s, Inches(1), Inches(1), Inches(5), Inches(3),
                                  ["A", "B", "C"], [10, 20, 30],
                                  [RED, BLUE, GREEN])
    assert cf and cf.chart


def test_chart_bar():
    _, s = make_blank()
    cf = chart.add_bar_chart(s, Inches(1), Inches(1), Inches(5), Inches(3),
                               ["A", "B"], [40, 60])
    assert cf


def test_chart_grouped_col():
    _, s = make_blank()
    cf = chart.add_grouped_col(s, Inches(1), Inches(1), Inches(6), Inches(3),
                                 ["Q1", "Q2"],
                                 {"A": ([1, 2], RED), "B": ([3, 4], BLUE)})
    assert cf


def test_chart_stacked_col():
    _, s = make_blank()
    cf = chart.add_stacked_col(s, Inches(1), Inches(1), Inches(6), Inches(3),
                                 ["Q1"],
                                 [("A", [10], RED), ("B", [20], BLUE)],
                                 percent=False)
    assert cf


def test_chart_stacked_bar():
    _, s = make_blank()
    cf = chart.add_stacked_bar(s, Inches(1), Inches(1), Inches(6), Inches(3),
                                 ["A", "B"],
                                 [("X", [50, 30], RED), ("Y", [50, 70], BLUE)],
                                 percent=True)
    assert cf


def test_chart_pie():
    _, s = make_blank()
    cf = chart.add_pie_chart(s, Inches(1), Inches(1), Inches(5), Inches(5),
                               ["A", "B", "C"], [40, 30, 30],
                               [RED, BLUE, GREEN])
    assert cf


def test_chart_doughnut():
    _, s = make_blank()
    cf = chart.add_doughnut_chart(s, Inches(1), Inches(1), Inches(5), Inches(5),
                                    ["A", "B"], [60, 40], [RED, BLUE])
    assert cf


def test_chart_line():
    _, s = make_blank()
    cf = chart.add_line_chart(s, Inches(1), Inches(1), Inches(6), Inches(3),
                                ["1月", "2月", "3月"],
                                {"GMV": ([100, 120, 140], RED)})
    assert cf


def test_chart_area():
    _, s = make_blank()
    cf = chart.add_area_chart(s, Inches(1), Inches(1), Inches(6), Inches(3),
                                ["1月", "2月"],
                                {"A": ([10, 20], RED), "B": ([15, 25], BLUE)},
                                stacked=True)
    assert cf


def test_chart_radar():
    _, s = make_blank()
    cf = chart.add_radar_chart(s, Inches(1), Inches(1), Inches(5), Inches(5),
                                 ["维度 1", "维度 2", "维度 3"],
                                 {"品牌 A": ([3, 4, 5], RED)})
    assert cf


def test_chart_scatter():
    _, s = make_blank()
    cf = chart.add_scatter_chart(
        s, Inches(1), Inches(1), Inches(6), Inches(4),
        {"组 A": [(1, 2), (2, 4), (3, 6)]},
        [BLUE],
    )
    assert cf


def test_chart_bubble():
    _, s = make_blank()
    cf = chart.add_bubble_chart(
        s, Inches(1), Inches(1), Inches(6), Inches(4),
        {"组 A": [(1, 2, 5), (2, 4, 10)]},
        [RED],
    )
    assert cf


def test_diagram_candlestick():
    """K 线图(走 diagrams 形状绘制,因为 python-pptx 不支持 STOCK_OHLC)。"""
    _, s = make_blank()
    diagrams.add_candlestick(
        s, Inches(1), Inches(1), Inches(8), Inches(4),
        ["1日", "2日", "3日", "4日"],
        [(100, 110, 95, 105),
         (105, 115, 100, 112),
         (112, 120, 108, 110),  # 跌
         (110, 118, 105, 115)],
    )
    assert len(s.shapes) > 0


# ══════════════════════════════════════════════════════════
# B 类:Diagram 测试(共 7 个)
# ══════════════════════════════════════════════════════════
def test_diagram_funnel():
    _, s = make_blank()
    diagrams.add_funnel(
        s, Inches(2), Inches(1), Inches(8), Inches(5),
        ["曝光", "点击", "下单"],
        [10000, 3000, 800],
    )
    assert len(s.shapes) > 0


def test_diagram_process_flow():
    _, s = make_blank()
    diagrams.add_process_flow(
        s, Inches(0.5), Inches(2), Inches(12), Inches(1),
        [("立势", "Q3"), ("冲刺", "Q4"), ("收割", "EOY")],
    )
    assert len(s.shapes) > 0


def test_diagram_sankey():
    _, s = make_blank()
    diagrams.add_sankey(
        s, Inches(2), Inches(1), Inches(8), Inches(5),
        [("源 A", 50), ("源 B", 50)],
        [("目标 X", 30), ("目标 Y", 70)],
        [(0, 0, 20), (0, 1, 30), (1, 1, 40), (1, 0, 10)],
    )
    assert len(s.shapes) > 0


def test_diagram_treemap():
    _, s = make_blank()
    diagrams.add_treemap(
        s, Inches(1), Inches(1), Inches(10), Inches(5),
        [("产品 A", 100), ("产品 B", 60), ("产品 C", 30)],
        [RED, BLUE, GREEN],
    )
    assert len(s.shapes) > 0


def test_diagram_sunburst():
    _, s = make_blank()
    diagrams.add_sunburst(
        s, Inches(2), Inches(1), Inches(8), Inches(5),
        {"分类 A": {"子 1": 10, "子 2": 20},
         "分类 B": {"子 3": 15}},
    )
    assert len(s.shapes) > 0


def test_diagram_gantt():
    _, s = make_blank()
    cf = diagrams.add_gantt(
        s, Inches(1), Inches(1), Inches(10), Inches(4),
        [("任务 A", 0, 5), ("任务 B", 3, 7), ("任务 C", 8, 4)],
    )
    assert cf is not None


def test_diagram_heatmap():
    _, s = make_blank()
    diagrams.add_heatmap(
        s, Inches(1), Inches(1), Inches(10), Inches(5),
        [[1, 2, 3], [4, 5, 6], [7, 8, 9]],
        ["X1", "X2", "X3"],
        ["Y1", "Y2", "Y3"],
    )
    assert len(s.shapes) > 0


# ══════════════════════════════════════════════════════════
# chart_picker 决策树测试
# ══════════════════════════════════════════════════════════
def test_picker_trend_recommends_line():
    """时间序列数据应推荐折线图。"""
    data = {
        "cats": ["1月", "2月", "3月", "4月"],
        "series": {"GMV": ([100, 120, 140, 130], None)},
    }
    recs = chart_picker.auto_select(data)
    assert recs[0]["type"] == "line"


def test_picker_composition_recommends_doughnut():
    """占比数据(总和接近 100)应推荐饼/环形。"""
    data = {
        "cats": ["A", "B", "C", "D"],
        "vals": [38, 23, 24, 15],  # 总和 100
    }
    recs = chart_picker.auto_select(data)
    # top 1 是 doughnut 或 pie
    assert recs[0]["type"] in ("doughnut", "pie")


def test_picker_flow_recommends_funnel():
    """流向意图应推荐漏斗或流程图。"""
    data = {"stages": ["A", "B", "C"]}
    recs = chart_picker.auto_select(data, intent="flow")
    assert recs[0]["type"] in ("funnel", "process_flow", "sankey")


def test_picker_explain_format():
    """explain 应该输出可读文本。"""
    recs = chart_picker.auto_select(
        {"cats": ["A", "B"], "vals": [50, 50]}
    )
    text = chart_picker.explain(recs)
    assert "推荐图表" in text
    assert "🟢" in text or "🟡" in text


# ══════════════════════════════════════════════════════════
# render_from_spec 端到端测试
# ══════════════════════════════════════════════════════════
def test_chart_render_from_spec_native():
    """chart.render_from_spec 应支持所有 A 类。"""
    _, s = make_blank()
    spec = {"type": "line", "cats": ["1", "2"],
            "series": {"X": ([1, 2], BLUE)}}
    cf = chart.render_from_spec(s, Inches(1), Inches(1),
                                 Inches(5), Inches(3), spec)
    assert cf


def test_diagrams_render_from_spec():
    """diagrams.render_from_spec 应支持所有 B 类。"""
    _, s = make_blank()
    spec = {"type": "funnel",
            "stages": ["X", "Y"], "values": [100, 50]}
    diagrams.render_from_spec(s, Inches(1), Inches(1),
                              Inches(8), Inches(4), spec)
    assert len(s.shapes) > 0


# ══════════════════════════════════════════════════════════
# chart_remake 工厂函数测试
# ══════════════════════════════════════════════════════════
def test_remake_kind_classification():
    """attach_chart_to_slide 应根据 type 自动判断 kind。"""
    from scripts.pipeline.chart_remake import (
        _kind_of, NATIVE_CHART_TYPES, DIAGRAM_TYPES,
    )
    for t in NATIVE_CHART_TYPES:
        assert _kind_of(t) == "native_chart"
    for t in DIAGRAM_TYPES:
        assert _kind_of(t) == "diagram"


if __name__ == "__main__":
    passed = 0
    failed = 0
    for name, fn in list(globals().items()):
        if name.startswith("test_") and callable(fn):
            try:
                fn()
                print(f"✅ {name}")
                passed += 1
            except Exception as e:
                print(f"❌ {name}: {type(e).__name__}: {e}")
                failed += 1
    print(f"\n{passed} passed, {failed} failed")
