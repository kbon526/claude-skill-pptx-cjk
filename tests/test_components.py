"""组件库冒烟测试 — 确保每个 add_* 函数能调用且返回非 None。"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.util import Inches

from scripts.components import text, bullets, kpi, table, chart, hero, accent


def make_blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def test_text_add_textbox():
    _, slide = make_blank_slide()
    tb = text.add_textbox(slide, Inches(1), Inches(1), Inches(4), Inches(1),
                           "测试中文")
    assert tb is not None
    assert tb.has_text_frame


def test_bullets_add_bullets():
    _, slide = make_blank_slide()
    bullets.add_bullets(slide, Inches(0.5), Inches(2), Inches(6),
                        ["要点 1", "要点 2", "要点 3"])
    # 红点 + 文字 = 6 个 shape
    assert len(slide.shapes) == 6


def test_kpi_add_kpi():
    _, slide = make_blank_slide()
    bg = kpi.add_kpi(slide, Inches(8), Inches(2), val="37%", label="SOV")
    assert bg is not None


def test_kpi_grid_2x2():
    _, slide = make_blank_slide()
    kpis = [
        {"val": "37%", "label": "SOV"},
        {"val": "5.2x", "label": "Reach"},
        {"val": "¥800万", "label": "Budget"},
        {"val": "12%", "label": "Share"},
    ]
    shapes = kpi.add_kpi_grid_2x2(slide, kpis, top_y=Inches(1.5))
    assert len(shapes) == 4


def test_table_add_table():
    _, slide = make_blank_slide()
    data = [
        ["品牌", "SOV", "增速"],
        ["品牌 A", "38%", "+5%"],
        ["品牌 B", "23%", "+12%"],
        ["AcmeCorp", "12%", "+15%"],
    ]
    ts = table.add_table(slide, Inches(1), Inches(2), Inches(8), Inches(2.5),
                         data)
    assert ts is not None


def test_chart_add_column_chart():
    _, slide = make_blank_slide()
    cf = chart.add_column_chart(
        slide, Inches(1), Inches(2), Inches(6), Inches(4),
        cats=["A", "B", "C"], vals=[10, 20, 30],
    )
    assert cf is not None
    assert cf.chart is not None


def test_chart_add_doughnut_chart():
    from scripts.core.color import RED, BLUE, GREEN
    _, slide = make_blank_slide()
    cf = chart.add_doughnut_chart(
        slide, Inches(1), Inches(2), Inches(5), Inches(5),
        cats=["X", "Y", "Z"], vals=[33, 33, 34],
        colors=[RED, BLUE, GREEN],
    )
    assert cf is not None


def test_chart_render_from_spec():
    _, slide = make_blank_slide()
    spec = {"type": "column", "cats": ["A", "B"], "vals": [1, 2]}
    cf = chart.render_from_spec(slide, Inches(1), Inches(1), Inches(4), Inches(3),
                                 spec)
    assert cf is not None


def test_hero_add_section_divider():
    prs = Presentation()
    slide = hero.add_section_divider(prs, 1, "测试章节", "副标题")
    assert slide is not None


def test_accent_add_red_label():
    _, slide = make_blank_slide()
    s = accent.add_red_label(slide, Inches(1), Inches(1),
                              Inches(2), Inches(0.4), "热门")
    assert s is not None


def test_accent_add_phase_track():
    from scripts.core.color import (
        PHASE_AWARENESS, PHASE_INTEREST, PHASE_CONVERSION
    )
    _, slide = make_blank_slide()
    accent.add_phase_track(slide, Inches(1.5), [
        ("立势期", "品牌建设", PHASE_AWARENESS),
        ("冲刺期", "兴趣获取", PHASE_INTEREST),
        ("收割期", "转化优化", PHASE_CONVERSION),
    ])
    # 3 个色块,每个色块本身是 1 个 shape
    assert len(slide.shapes) >= 3


def test_registry_init_and_add():
    from scripts.core.registry import (
        init_registry, add_slide, attach_visual,
        validate, summary, KIND_NATIVE_CHART,
    )
    reg = init_registry("Test Deck", brand="Test", version="v1")
    assert reg["meta"]["deck_title"] == "Test Deck"

    add_slide(reg, "slide_01", title="测试页", layout="default")
    attach_visual(reg, "slide_01", kind=KIND_NATIVE_CHART,
                  spec={"type": "column", "cats": ["A"], "vals": [1]})

    errors = validate(reg)
    assert errors == [], f"Registry validation failed: {errors}"

    s = summary(reg)
    assert "Test Deck" in s
    assert "slide_01" in s


if __name__ == "__main__":
    for name, fn in list(globals().items()):
        if name.startswith("test_") and callable(fn):
            try:
                fn()
                print(f"✅ {name}")
            except Exception as e:
                print(f"❌ {name}: {type(e).__name__}: {e}")
